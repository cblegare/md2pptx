#!/usr/bin/env python3

"""
md2pptx - Converts (a subset of) Markdown to Powerpoint (PPTX)

First argument is file to write to

Reads from stdin
"""

import re
import sys
import os
import csv
import time
import collections
import collections.abc
from pptx import Presentation
from pptx import __version__ as pptx_version
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor, MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER, PP_MEDIA_TYPE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.action import PP_ACTION
from pptx.enum.dml import MSO_PATTERN

import struct
import datetime
import html.parser
from pptx.oxml.xmlchemy import OxmlElement
from pathlib import Path
import urllib.request
import tempfile
import copy
import platform
import shutil
import socket
from pptx.oxml import parse_xml
import uuid

import md2pptx.funnel
import md2pptx.runPython
from md2pptx.card import Card
from md2pptx.rectangle import Rectangle
from md2pptx.colour import *
from md2pptx.paragraph import *
from md2pptx.symbols import resolveSymbols
import md2pptx.globals
from md2pptx.processingOptions import *


from lxml import etree
from lxml.html import fromstring

# Try to import CairoSVG - which might not be installed.
# Flag availability or otherwise
try:
    import cairosvg
    from cairosvg import helpers

    have_cairosvg = True

except:
    have_cairosvg = False

# Try to import Pillow - which might not be installed.
# Flag availability or otherwise
try:
    import PIL

    have_pillow = True

except:
    have_pillow = False

# Try to import graphviz - which might not be installed.
# Flag availability or otherwise

try:
    import graphviz

    have_graphviz = True

except:
    have_graphviz = False


md2pptx_level = "5.4"
md2pptx_date = "23 February, 2025"

namespaceURL = {
    "mc":  "http://schemas.openxmlformats.org/markup-compatibility/2006",
    "p":   "http://schemas.openxmlformats.org/presentationml/2006/main",
    "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
    "p15": "http://schemas.microsoft.com/office/powerpoint/2012/main",
    "a":   "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r":   "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
}

def namespacesFragment(prefixes):
    xml = ""
    for prefix in prefixes:
        xml += 'xmlns:' + prefix + '="' + namespaceURL[prefix] +'" '

    return xml

class SlideInfo:
    def __init__(
        self,
        titleText,
        subtitleText,
        blockType,
        bullets,
        tableRows,
        cards,
        code,
        sequence,
    ):
        self.titleText = titleText
        self.subtitleText = subtitleText
        self.blockType = blockType
        self.bullets = bullets
        self.tableRows = tableRows
        self.cards = cards
        self.code = code
        self.sequence = sequence


# Information about a single table. (A slide might have more than one - or none.)
class TableInfo:
    def __init__(self, tableRows, tableCaption):
        self.tableRows = []
        self.tableCaption = ""


# Information about a video
class AudioVideoInfo:
    def __init__(self, elementString):
        audioVideoElement = fromstring(elementString)

        if "src" in audioVideoElement.attrib.keys():
            self.source = audioVideoElement.attrib["src"]
        else:
            self.source = ""

        if audioVideoElement.tag == "audio":
            # audio element doesn't have width or height attributes so make it square
            self.width = 1024
            self.height = 1024
        else:
            # video element can have width and height attributes
            # Default is 4:3
            if "width" in audioVideoElement.attrib.keys():
                self.width = int(audioVideoElement.attrib["width"])
            else:
                self.width = 1024

            if "height" in audioVideoElement.attrib.keys():
                self.height = int(audioVideoElement.attrib["height"])
            else:
                self.height = 768

        self.aspectRatio = self.width / self.height

        if "poster" in audioVideoElement.attrib.keys():
            self.poster = audioVideoElement.attrib["poster"]
        else:
            self.poster = None



# Get a picture's rId
def get_picture_rId(picture):
    rId = picture._element.xpath("./p:blipFill/a:blip/@r:embed")[0]
    return rId

# Adds a picture as a background
def add_background(presentation, slide, picture):
    # Add the picture with zero dimensions
    picture = slide.shapes.add_picture(picture,0,0,0,00)

    # Get the RId for this tiny picture = as we'll need that to set the background
    rId = get_picture_rId(picture)

    # find the cSld element to attach the XML to
    cSld =slide._element.xpath("./p:cSld")[0]

    # Remove any pre-existing bg element
    bg =slide._element.xpath("./p:bg")
    if bg != []:
        cSld.remove(bg)


    # Confect the XML
    xml = ""
    xml += '    <p:bg ' + namespacesFragment(["a","p","r"]) + '>\n'
    xml += '       <p:bgPr>\n'
    xml += '         <a:blipFill xmlns:a="' + namespaceURL["a"] + '" dpi="0" rotWithShape="1">\n'
    xml += '           <a:blip r:embed="' + rId +'">\n'
    xml += '             <a:lum />\n'
    xml += '               </a:blip>\n'
    xml += '           <a:srcRect />\n'
    xml += '           <a:stretch>\n'
    xml += '             <a:fillRect t="0" b="0" />\n'
    xml += '          </a:stretch>\n'
    xml += '        </a:blipFill>\n'
    xml += '        <a:effectLst/>\n'
    xml += '      </p:bgPr>\n'
    xml += '    </p:bg>\n'

    # Parse this XML
    parsed_xml = parse_xml(xml)

    # Insert the parsed XML fragment as a child of the cSld element
    cSld.insert(0, parsed_xml)

    # Delete the original picture
    deleteSimpleShape(picture)

# Find the extLst element - if it exists in presentation.xml
def findExtLst(prs):
    for child in prs._element.getchildren():
        if child.tag.endswith("}extLst"):
            return child

    return None


def addSlide(presentation, slideLayout, slideInfo=None):
    slide = presentation.slides.add_slide(slideLayout)
    slide.slideInfo = slideInfo

    backgroundImage = md2pptx.globals.processingOptions.getCurrentOption("backgroundImage")

    if backgroundImage != "":
        add_background(presentation, slide, backgroundImage)
    return slide


def createSectionsXML(prs):
    sectionSlideLayout = md2pptx.globals.processingOptions.getCurrentOption("SectionSlideLayout")

    xml = '  <p:ext ' + namespacesFragment(["p"])

    # ext URI has to be {521415D9-36F7-43E2-AB2F-B90AF26B5E84} as it's a registered extension
    xml += '    uri="{521415D9-36F7-43E2-AB2F-B90AF26B5E84}">\n'

    xml += '    <p14:sectionLst ' + namespacesFragment(["p14"]) + '>\n'

    sectionCount = 0
    for slide in prs.slides:
        slideID = str(slide.slide_id)

        for idx, slide_layout in enumerate(prs.slide_layouts):
            if slide.slide_layout == slide_layout:
                layoutNumber = idx
                break

        if layoutNumber == sectionSlideLayout:
            # Have a section to contribute
            sectionCount += 1

            # Confect section name from first part of section slide title
            title = findTitleShape(slide)
            sectionName = title.text.split("\n")[0]

            # Clean up section name
            sectionName = "".join(
                letter
                for letter in sectionName
                if (
                    (letter.isalnum())
                    | (letter in "&-+")
                    | (letter in "!/*")
                    | (letter == " ")
                )
            )

            sectionName = (
                sectionName.replace("& ", "&amp; ")
                .replace("\r", " ")
                .replace("\n", " ")
            )

            # section URI's just need to be a GUID wrapped in braces
            xml += (
                '      <p14:section name="'
                + sectionName
                + '" id="{'
                + str(uuid.uuid4()).upper()
                + '}">\n'
            )

            # Only the first slide in the section is added - as section will continue until the next section
            # anyway
            xml += "        <p14:sldIdLst>\n"
            xml += '          <p14:sldId id="' + slideID + '" />\n'
            xml += "        </p14:sldIdLst>\n"
            xml += "      </p14:section>\n"

    # Close out the section list
    xml += "    </p14:sectionLst>\n"

    # Close out the sections extension
    xml += "  </p:ext>\n"

    parsed_xml = parse_xml(xml)

    return parsed_xml, sectionCount


def createExpandingSections(prs):
    # Use the slides' layouts to create an XML fragment with sections in
    xmlFragment, sectionCount = createSectionsXML(prs)

    if sectionCount > 0:
        # Have sections to insert as an XML fragment
        if (extLstElement := findExtLst(prs)) is not None:
            # Need to remove the extension list element before adding a new one
            prs._element.remove(extLstElement)

        # Insert a new extension list element
        extLst = OxmlElement("p:extLst")
        prs._element.insert(-1, extLst)

        # Insert the fragment in the extension list in presentation.xml
        extLst.insert(0, xmlFragment)


def deleteSlide(prs, slideNumber):
    rId = prs.slides._sldIdLst[slideNumber].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[slideNumber]


def startswithOneOf(haystack, needleList):
    for needle in needleList:
        if haystack.startswith(needle):
            return True

    return False


# Splits a string into words, converting each word to an integer. Returns them as a
# sorted list
def sortedNumericList(string):
    return sorted(list(map(int, set(string.split()))))


def substituteFooterVariables(footerText, liveFooters):
    # Decide if the footer should be a live link to the section slide
    wantLiveFooter = (
        (prs.lastSectionSlide is not None)
        & (footerText.find("<section") > -1)
        & (liveFooters == "yes")
    )

    # Substitute any section title occurrences
    sectionTitleLines = resolveSymbols(prs.lastSectionTitle).split("<br/>")

    footerText = footerText.replace("<section>", sectionTitleLines[0])
    footerText = footerText.replace("<section1>", sectionTitleLines[0])

    if len(sectionTitleLines) > 1:
        footerText = footerText.replace("<section2>", sectionTitleLines[1])

    if len(sectionTitleLines) > 2:
        footerText = footerText.replace("<section3>", sectionTitleLines[2])

    # Substitute any presentation title occurrences
    presTitleLines = resolveSymbols(prs.lastPresTitle).split("<br/>")

    footerText = footerText.replace("<presTitle>", presTitleLines[0])
    footerText = footerText.replace("<presTitle1>", presTitleLines[0])

    if len(presTitleLines) > 1:
        footerText = footerText.replace("<presTitle2>", presTitleLines[1])

    if len(presTitleLines) > 2:
        footerText = footerText.replace("<presTitle3>", presTitleLines[2])

    # Substitute any presentation subtitle occurrences
    presSubtitleLines = resolveSymbols(prs.lastPresSubtitle).split("<br/>")

    footerText = footerText.replace("<presSubtitle>", presSubtitleLines[0])
    footerText = footerText.replace("<presSubtitle1>", presSubtitleLines[0])

    if len(presSubtitleLines) > 1:
        footerText = footerText.replace("<presSubtitle2>", presSubtitleLines[1])

    if len(presSubtitleLines) > 2:
        footerText = footerText.replace("<presSubtitle3>", presSubtitleLines[2])

    # Make newlines happen
    footerText = footerText.replace("<br/>", "\n")

    return footerText, wantLiveFooter

def addTableShadow(t):
    tblPr = t._tbl.getchildren()[0]

    xml = """
    <a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <a:outerShdw xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" blurRad="50800" dist="190500" dir="2700000" algn="tl" rotWithShape="0">
            <a:prstClr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="black">
                <a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="40000" />
            </a:prstClr>
        </a:outerShdw>
    </a:effectLst>
"""
    # Parse this XML
    parsed_xml = parse_xml(xml)

    # Insert the parsed XML fragment as a child of the pPr element
    tblPr.insert(0, parsed_xml)

def _applyCellBorderStyling(
    tcPr, linePosition, lineWidthMultiplier=1, lineCount=1, lineColour="000000"
):
    # How wide, relatively speaking to make the lines
    lineWidth = int(12700 * lineWidthMultiplier)

    # Whether the line should be single or double
    if lineCount == 2:
        lineCountValue = "dbl"
    else:
        lineCountValue = "sng"

    if linePosition == "l":
        elementName = "a:lnL"
    elif linePosition == "r":
        elementName = "a:lnR"
    elif linePosition == "t":
        elementName = "a:lnT"
    else:
        elementName = "a:lnB"

    lnX = OxmlElement(elementName)

    lnX.attrib.update(
        {"w": str(lineWidth), "cap": "flat", "cmpd": lineCountValue, "algn": "ctr"}
    )

    solidFill = OxmlElement("a:solidFill")
    srgbClr = OxmlElement("a:srgbClr")
    srgbClr.attrib.update({"val": lineColour})

    solidFill.append(srgbClr)
    lnX.append(solidFill)

    tcPr.append(lnX)


def applyCellBorderStyling(
    cell, cellBorderStyling, lineWidthMultiplier, lineCount, lineColour
):
    if cellBorderStyling == "":
        # No cell border styling required
        return

    # Get any existing cell properties element - or make one
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    # Draw any cell borders. More than one might apply
    if cellBorderStyling.find("l") > -1:
        _applyCellBorderStyling(tcPr, "l", lineWidthMultiplier, lineCount, lineColour)
    if cellBorderStyling.find("r") > -1:
        _applyCellBorderStyling(tcPr, "r", lineWidthMultiplier, lineCount, lineColour)
    if cellBorderStyling.find("t") > -1:
        _applyCellBorderStyling(tcPr, "t", lineWidthMultiplier, lineCount, lineColour)
    if cellBorderStyling.find("b") > -1:
        _applyCellBorderStyling(tcPr, "b", lineWidthMultiplier, lineCount, lineColour)


# Apply table line styling
def applyTableLineStyling(
    table,
    processingOptions,
):
    wholeTableLineStyling = processingOptions.getCurrentOption("addTableLines")
    linedColumns = processingOptions.getCurrentOption("addTableColumnLines")
    linedRows = processingOptions.getCurrentOption("addTableRowLines")

    lastRow = len(table.rows) - 1

    # Create blank cell styling matrix
    cellStyling = []
    for rowNumber, row in enumerate(table.rows):
        rowStyling = []
        for cell in row.cells:
            rowStyling.append("")
        cellStyling.append(rowStyling)

    # apply any "whole table" styling - from addTableLines
    if wholeTableLineStyling == "box":
        # Line around the table
        for rowNumber, row in enumerate(table.rows):
            # Figure out whether row is top, middle, or bottom
            if rowNumber == 0:
                rowStyling = "t"
            elif rowNumber == lastRow:
                rowStyling = "b"
            else:
                rowStyling = ""

            lastColumn = len(row.cells) - 1

            for columnNumber, cell in enumerate(row.cells):
                if columnNumber == 0:
                    columnStyling = "l"
                elif columnNumber == lastColumn:
                    columnStyling = "r"
                else:
                    columnStyling = ""
                cellStyling[rowNumber][columnNumber] = rowStyling + columnStyling

    elif wholeTableLineStyling == "all":
        # All edges of all cells have lines
        for rowNumber, row in enumerate(table.rows):
            lastColumn = len(row.cells) - 1

            for columnNumber, cell in enumerate(row.cells):
                cellStyling[rowNumber][columnNumber] = "tlbr"

    # Apply any row styling - from addTableColumnLines
    for rowNumber, row in enumerate(table.rows):
        if rowNumber + 1 in linedRows:
            # Line after this row so below
            for columnNumber, cell in enumerate(row.cells):
                cellStyling[rowNumber][columnNumber] = (
                    cellStyling[rowNumber][columnNumber] + "b"
                )

        elif rowNumber in linedRows:
            # Line before this row so above
            for columnNumber, cell in enumerate(row.cells):
                cellStyling[rowNumber][columnNumber] = (
                    cellStyling[rowNumber][columnNumber] + "t"
                )

    # Apply any column styling - from addTableRowLines
    for rowNumber, row in enumerate(table.rows):
        for columnNumber, cell in enumerate(row.cells):
            if columnNumber + 1 in linedColumns:
                # Line after this column so to right
                cellStyling[rowNumber][columnNumber] = (
                    cellStyling[rowNumber][columnNumber] + "r"
                )

            elif columnNumber + 1 in linedColumns:
                # Line after this column so to left
                cellStyling[rowNumber][columnNumber] = (
                    cellStyling[rowNumber][columnNumber] + "r"
                )

    # Apply the styling from the matrix to all cells
    for rowNumber, row in enumerate(table.rows):
        for columnNumber, cell in enumerate(row.cells):
            applyCellBorderStyling(
                cell,
                cellStyling[rowNumber][columnNumber],
                md2pptx.globals.processingOptions.getCurrentOption("addTableLineWidth"),
                md2pptx.globals.processingOptions.getCurrentOption("addTableLineCount"),
                md2pptx.globals.processingOptions.getCurrentOption("addTableLineColour"),
            )



def reportSlideTitle(slideNumber, indent, titleText):
    print(str(slideNumber).rjust(4) + " " + ("  " * indent) + titleText)


def reportGraphicFilenames(leftFilename, rightFilename=""):
    if rightFilename == "":
        print("             ---> " + leftFilename.ljust(30))
    else:
        print("             ---> " + leftFilename.ljust(30) + " , " + rightFilename)


# Given current indenting regime calculate what level the bullet / number is at
def calculateIndentationLevel(firstNonSpace, indentSpaces):
    return int(firstNonSpace / indentSpaces)


# Calculate picture dimensions given its natural height and bounds
def scalePicture(maxPicWidth, maxPicHeight, imageWidth, imageHeight):
    heightIfWidthUsed = maxPicWidth * imageHeight / imageWidth
    widthIfHeightUsed = maxPicHeight * imageWidth / imageHeight

    if heightIfWidthUsed > maxPicHeight:
        # Use the height to scale
        usingHeightToScale = True

        picWidth = widthIfHeightUsed
        picHeight = maxPicHeight

    else:
        # Use the width to scale
        usingHeightToScale = False

        picWidth = maxPicWidth
        picHeight = heightIfWidthUsed
    return (picWidth, picHeight, usingHeightToScale)


def parseMedia(cellString, graphicCount):
    graphicTitle = ""
    HTML = ""
    audioVideoInfo = None
    graphicHref = ""
    GraphicFilename = ""
    printableGraphicFilename = ""

    graphicCount += 1

    if videoRegexMatch := videoRegex.match(cellString):
        # Cell contains a video
        audioVideoInfo = AudioVideoInfo(cellString)
        _, printableGraphicFilename = handleWhateverGraphicType(audioVideoInfo.source)

    elif audioRegexMatch := audioRegex.match(cellString):
        # Cell contains an audio
        audioVideoInfo = AudioVideoInfo(cellString)
        _, printableGraphicFilename = handleWhateverGraphicType(audioVideoInfo.source)

    elif clickableGraphicMatch := clickableGraphicRegex.match(cellString):
        # Cell contains a clickable graphic
        graphicTitle = clickableGraphicMatch.group(1)
        GraphicFilename = clickableGraphicMatch.group(2)
        graphicHref = clickableGraphicMatch.group(3)

        (
            GraphicFilename,
            printableGraphicFilename,
        ) = handleWhateverGraphicType(GraphicFilename)

    elif graphicMatch := graphicRegex.match(cellString):
        # Cell contains a non-clickable graphic
        graphicTitle = graphicMatch.group(1)
        GraphicFilename = graphicMatch.group(2)

        (
            GraphicFilename,
            printableGraphicFilename,
        ) = handleWhateverGraphicType(GraphicFilename)

    else:
        # Not a graphic or video
        GraphicFilename = ""
        printableGraphicFilename = ""
        HTML = cellString
        graphicCount -= 1

    return (
        graphicTitle,
        GraphicFilename,
        printableGraphicFilename,
        graphicHref,
        HTML,
        audioVideoInfo,
        graphicCount,
    )


# Send a shape to the back on a slide
def sendToBack(shapes, shape):
    firstShapeElement = shapes[0]._element
    firstShapeElement.addprevious(shape._element)


# Turn a paragraph into a numbered inList item
def makeNumberedListItem(p):
    if (
        p._element.getchildren()[0].tag
        == "{http://schemas.openxmlformats.org/drawingml/2006/main}pPr"
    ):
        pPr = p._element.getchildren()[0]
        if len(pPr.getchildren()) > 0:
            # Remove Default Text Run Properties element - if present
            x = pPr.getchildren()[0]
            if x.tag == "{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr":
                pPr.remove(x)
    else:
        pPr = OxmlElement("a:pPr")
        p._element.insert(0, pPr)

    buFont = OxmlElement("a:buFont")
    buFont.set("typeface", "+mj-lt")
    pPr.append(buFont)

    buAutoNum = OxmlElement("a:buAutoNum")
    buAutoNum.set("type", "arabicPeriod")
    pPr.append(buAutoNum)

# Add a drop shadow to a shape
def createShadow(shape):
    if "Table" in shape.__class__.__name__:
        # Table has to be handled differently
        return addTableShadow(shape)

    el = OxmlElement("a:effectLst")

    spPr = shape.fill._xPr

    spPr.append(el)

    outerShdw = OxmlElement("a:outerShdw")
    outerShdw.set("algn", "tl")
    outerShdw.set("blurRad", "50800")
    outerShdw.set("dir", "2700000")
    outerShdw.set("dist", "95250")
    outerShdw.set("rotWithShape", "0")

    el.append(outerShdw)

    prstClr = OxmlElement("a:prstClr")
    prstClr.set("val", "black")

    outerShdw.append(prstClr)

    alpha = OxmlElement("a:alpha")
    alpha.set("val", "40000")

    prstClr.append(alpha)


# Clone a shape in a slide and return the new shape.
# (This is a deep copy so the new shape will have the same
# eg bullet style as the source shape)
def addClonedShape(slide, shape1):
    # Clone the element for the shape
    el1 = shape1.element
    el2 = copy.deepcopy(el1)

    # Insert the cloned element into the shape tree
    slide.shapes._spTree.insert_element_before(el2, "p:extLst")

    # Return the shape associated with this new element
    return slide.shapes[-1]


# Following functions are workarounds for python-pptx not having these functions for the font object
def set_subscript(font):
    if font.size is None:
        font._element.set("baseline", "-50000")
        return

    if font.size < Pt(24):
        font._element.set("baseline", "-50000")
    else:
        font._element.set("baseline", "-25000")


def set_superscript(font):
    if font.size is None:
        font._element.set("baseline", "60000")
        return

    if font.size < Pt(24):
        font._element.set("baseline", "60000")
    else:
        font._element.set("baseline", "30000")


def setStrikethrough(font):
    font._element.set("strike", "sngStrike")


def setHighlight(run, color):
    # get run properties
    rPr = run._r.get_or_add_rPr()

    # Create highlight element
    hl = OxmlElement("a:highlight")

    # Create specify RGB Colour element with color specified
    srgbClr = OxmlElement("a:srgbClr")
    setattr(srgbClr, "val", color)

    # Add colour specification to highlight element
    hl.append(srgbClr)

    # Add highlight element to run properties
    rPr.append(hl)

    return run


# Get the slide object the run is in
def SlideFromRun(run):
    return run._parent._parent._parent._parent._parent


# Get the slide object the picture is in
def SlideFromPicture(picture):
    return picture._parent._parent


# Creates a hyperlink to another slide and/or a tooltip - for a
# text run
# Note: To get just a tooltip make to_slide be the source slide
#       so it links to itself.
def createRunHyperlinkOrTooltip(run, to_slide, tooltipText=""):
    # Get hold of the shape the run is in
    if run._parent._parent._parent.__class__.__name__ == "_Cell":
        # Run in a table cell has to be handled differently
        shape = (
            run._parent._parent._parent._parent._parent._parent._parent._graphic_frame
        )
    else:
        # Ordinary text run
        shape = run._parent._parent._parent

    if to_slide == None:
        to_slide = SlideFromRun(run)
    hl = run.hyperlink
    sca = shape.click_action
    sca_hl = sca.hyperlink

    # Add a click action to generate an internal hyperlink address
    sca.target_slide = to_slide

    # Use that internal hyperlink address for the run
    hl.address = sca_hl.address

    # Also clone the hyperlink click action
    hl._hlinkClick.action = sca_hl._hlink.action

    if tooltipText != "":
        hl._hlinkClick.set("tooltip", tooltipText)

    # Also clone the hyperlink rId
    hl._hlinkClick.rId = sca_hl._hlink.rId

    # Delete the shape click action
    sca.target_slide = None


# Creates a hyperlink to another slide or a URL and/or a tooltip - for a
# picture
# Note: To get just a tooltip make to_slide be the source slide
#       so it links to itself.
def createPictureHyperlinkOrTooltip(picture, target, tooltipText=""):
    if target == None:
        # If neither a tooltip nor a target slide then return having
        # done nothing
        if tooltipText == "":
            return

        # Tooltip but no target slide
        target = SlideFromPicture(picture)
        picture.click_action.target_slide = target
    elif target.__class__.__name__ == "str":
        # Is a URL
        picture.click_action.hyperlink.address = target
        # URL might be a macro reference
        if target[:11] == "ppaction://":
            # URL is indeed a macro reference, so treat it as such
            picture.click_action.hyperlink._hlink.set("action", target)
    else:
        # Target is a slide
        picture.click_action.target_slide = target

    if tooltipText != "":
        picture.click_action.hyperlink._hlink.set("tooltip", tooltipText)


# If a tooltip has been set return it else return an empty string
def getPictureTooltip(picture):
    if picture.click_action.hyperlink._hlink != None:
        # There is a tooltip
        return picture.click_action.hyperlink._hlink.get("tooltip")
    else:
        # There is no tooltip
        return ""


# Create hyperlink and optional tooltip from a shape eg Chevron
def createShapeHyperlinkAndTooltip(shape, to_slide, tooltipText=""):
    shape.click_action.target_slide = to_slide
    hl = shape.click_action.hyperlink
    hl._hlink.set("tooltip", tooltipText)


def getGraphicDimensions(fname):
    """Determine the image type of fhandle and return its size.
    from draco"""
    try:
        with open(fname, "rb") as fhandle:
            head = fhandle.read(24)
            if len(head) != 24:
                return -1, -1
            fname2 = fname.lower()
            if fname2.endswith(".png"):
                check = struct.unpack(">i", head[4:8])[0]
                if check != 0x0D0A1A0A:
                    return -1, -1
                width, height = struct.unpack(">ii", head[16:24])
            elif fname2.endswith(".gif"):
                width, height = struct.unpack("<HH", head[6:10])
            elif fname.endswith((".jpeg", ".jpg")):
                try:
                    fhandle.seek(0)  # Read 0xff next
                    size = 2
                    ftype = 0
                    while not 0xC0 <= ftype <= 0xCF:
                        fhandle.seek(size, 1)
                        byte = fhandle.read(1)
                        while ord(byte) == 0xFF:
                            byte = fhandle.read(1)
                        ftype = ord(byte)
                        size = struct.unpack(">H", fhandle.read(2))[0] - 2
                    # We are at a SOFn block
                    fhandle.seek(1, 1)  # Skip 'precision' byte.
                    height, width = struct.unpack(">HH", fhandle.read(4))
                except Exception:  # IGNORE:W0703
                    return
            else:
                return -1, -1

            return width, height

    except EnvironmentError:
        return -1, -1


def getVideoInfo(audioVideoInfo):
    if audioVideoInfo.source.find("://") > -1:
        # Video would be sourced from the web
        try:
            operUrl = urllib.request.urlopen(audioVideoInfo.source)
        except urllib.error.HTTPError as e:
            return -1, -1, "Web", None

        except socket.error as s:
            return -1, -1, "Web", None

        data = operUrl.read()

        return audioVideoInfo.width, audioVideoInfo.height, "Web", data
    else:
        # Video would be sourced from a local file
        try:
            fhandle = open(audioVideoInfo.source, "rb")
        except EnvironmentError:
            return -1, -1, "Local", None

        return audioVideoInfo.width, audioVideoInfo.height, "Local", None


# Render a list of bullets
def renderText(shape, bullets):
    baseTextDecrement = md2pptx.globals.processingOptions.getCurrentOption("baseTextDecrement")
    baseTextSize = md2pptx.globals.processingOptions.getCurrentOption("baseTextSize")

    tf = shape.text_frame

    for bulletNumber, bullet in enumerate(bullets):
        para0 = tf.paragraphs[0]

        if bulletNumber == 0:
            # Don't need to create paragraph
            p = para0
        else:
            # We need a new paragraph
            p = tf.add_paragraph()

        # Set the paragraph's level - zero-indexed
        p.level = int(bullet[0])

        # Set the paragraph's font size, adjusted for level, if necessary
        if baseTextSize > 0:
            p.font.size = Pt(baseTextSize - p.level * baseTextDecrement)

        addFormattedText(p, bullet[1])

        if bullet[2] == "numbered":
            makeNumberedListItem(p)

    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def findTitleShape(slide):
    if slide.shapes.title == None:
        # Have to use first shape as title
        return slide.shapes[0]
    else:
        return slide.shapes.title


def findBodyShape(slide):
    if len(slide.shapes) > 1:
        return slide.shapes[1]
    elif slide.shapes.title == None:
        return slide.shapes[0]
    else:
        return None


# Returns a top, left, width, height for content to be rendered into
def getContentRect(presentation, slide, topOfContent, margin):
    numbersHeight = md2pptx.globals.processingOptions.getCurrentOption("numbersHeight")
    # Left and right are always defined by the margins
    rectLeft = margin
    rectWidth = presentation.slide_width - 2 * margin
    if topOfContent == 0:
        # There is no title on this slide
        rectTop = margin
        rectHeight = presentation.slide_height - margin - max(margin, numbersHeight)
    else:
        # There is a title on this slide
        rectTop = topOfContent + margin
        rectHeight = presentation.slide_height - rectTop - max(margin, numbersHeight)

    return (rectLeft, rectWidth, rectTop, rectHeight)


# Finds the title and adds the text to it, returning title bottom, title shape, and
# flattened title
def formatTitle(presentation, slide, titleText, titleFontSize, subtitleFontSize):
    marginBase = md2pptx.globals.processingOptions.getCurrentOption("marginBase")
    pageTitleAlign = md2pptx.globals.processingOptions.getCurrentOption("pagetitlealign")

    # Convert page title alignment text value to constant
    if pageTitleAlign == "left":
        titleAlignment = PP_ALIGN.LEFT
    elif pageTitleAlign == "right":
        titleAlignment = PP_ALIGN.RIGHT
    else:
        titleAlignment = PP_ALIGN.CENTER

    # Find title
    title = findTitleShape(slide)

    if titleText == "&nbsp;":
        deleteSimpleShape(title)

        return (marginBase, None, "<No title>")

    if md2pptx.globals.processingOptions.getCurrentOption("adjustTitles"):
        title.top = marginBase
        title.left = marginBase
        title.width = presentation.slide_width - marginBase * 2

    # Figure out how many lines title will need (ignoring overflow)
    titleLineCount = len(titleLines := titleText.split("<br/>"))

    # This will hold the flattened title lines to be printed
    flattenedTitleLines = []

    # Add the first line of the title text to the first paragraph
    firstTitleParagraph = title.text_frame.paragraphs[0]
    flattenedTitleLines.append(addFormattedText(firstTitleParagraph, titleLines[0]))

    # Set anchor to top
    title.text_frame.vertical_anchor = MSO_ANCHOR.TOP

    # Set this paragraph's font size using pageTitleSize
    firstTitleParagraph.font.size = Pt(titleFontSize)

    # No space before the paragraph
    firstTitleParagraph.space_after = Pt(0)

    # Set first title paragraph's alignment
    firstTitleParagraph.alignment = titleAlignment

    if subtitleFontSize == "same":
        subtitleFontSize = titleFontSize

    # If there are additional title lines then add them
    for lineNumber, titleLine in enumerate(titleLines[1:]):
        # Each title line requires a new paragraph
        newPara = title.text_frame.add_paragraph()

        # Use this new paragraph, adding flattened title line to the list
        flattenedTitleLines.append(addFormattedText(newPara, titleLine))

        # No space before the paragraph
        newPara.space_before = Pt(0)

        # Set this paragraph's font size using pageSubtitleSize
        newPara.font.size = Pt(subtitleFontSize)

        # Set this paragraph's alignment
        newPara.alignment = titleAlignment

    # Note: Working off pageTitleSize and pageSubtitleSize
    if md2pptx.globals.processingOptions.getCurrentOption("adjustTitles"):
        title.height = Pt(titleFontSize) + Pt(subtitleFontSize) * (titleLineCount - 1)

    # Massage title line for printing a little
    if titleLineCount > 1:
        flattenedTitleText = flattenedTitleLines[0] + " ..."
    else:
        flattenedTitleText = flattenedTitleLines[0]

    # Return where the next shape below the title would be - vertically
    return (title.top + title.height + Inches(0.1), title, flattenedTitleText)


# Parse the string after the e.g. ### for a displayable title and
# an optional heading reference
def parseTitleText(titleLineString):
    # Get rid of any cruft on the line
    slideTitleWithPossibleHref = titleLineString.strip().rstrip("#").rstrip()

    if hrefMatch := slideHrefRegex.match(slideTitleWithPossibleHref):
        # Use the explicit href
        slideTitle = hrefMatch.group(1)
        href = hrefMatch.group(2)
    else:
        # No href
        href = ""
        slideTitle = slideTitleWithPossibleHref

    return slideTitle, href



def addFooter(presentation, slideNumber, slide):
    numbersHeight = md2pptx.globals.processingOptions.getCurrentOption("numbersHeight")

    numbersglobals.fontsizespec = md2pptx.globals.processingOptions.getCurrentOption("numbersFontSize")
    if numbersglobals.fontsizespec == "":
        numbersFontSize = Pt(12)
    else:
        numbersFontSize = Pt(numbersglobals.fontsizespec)

    shapes = slide.shapes
    footer = shapes.add_textbox(
        Inches(0.1),
        presentation.slide_height - numbersHeight,
        Inches(0.2),
        numbersHeight / 2,
    )
    frame = footer.text_frame
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = str(slideNumber)
    font = run.font
    font.size = numbersFontSize


# Called "Simple" because more complex shapes might not work
def deleteSimpleShape(shape):
    if shape == None:
        return
    shapeElement = shape.element
    shapeElement.getparent().remove(shapeElement)


def createProcessingSummarySlide(presentation, rawMetadata):
    tableMargin = md2pptx.globals.processingOptions.getCurrentOption("tableMargin")
    pageTitleSize = md2pptx.globals.processingOptions.getCurrentOption("pageTitleSize")
    pageSubtitleSize = md2pptx.globals.processingOptions.getCurrentOption("pageSubtitleSize")
    # Use the first slide in the template presentation as the base
    slide = presentation.slides[0]

    # Delete any body shape - other than action buttons
    bodyShape = findBodyShape(slide)
    if bodyShape.name.startswith("Action Button:") is False:
        deleteSimpleShape(bodyShape)

    # Build "run time" text
    now = datetime.datetime.now()
    runTime = now.strftime("%H:%M").lstrip()
    runDate = now.strftime("%e %B, %G").lstrip()
    runDateTime = "Presentation built: " + runTime + " on " + runDate

    # Format title and add title text
    slideTitleBottom, title, flattenedTitle = formatTitle(
        presentation,
        slide,
        banner + "<br/>" + runDateTime,
        pageTitleSize,
        pageSubtitleSize,
    )

    # Work out how many pairs of columns we need
    if md2pptx.globals.processingOptions.hideMetadataStyle:
        # Adjust metadata item count to remove style.
        metadata = []
        for metadataItem in rawMetadata:
            if metadataItem[0].startswith("style.") == False:
                metadata.append(metadataItem)
    else:
        metadata = rawMetadata

    metadataRows = len(metadata)

    maxMetadataRowsPerColumn = 15
    if metadataRows > 4 * maxMetadataRowsPerColumn:
        metadataColumnPairs = 5
    elif metadataRows > 3 * maxMetadataRowsPerColumn:
        metadataColumnPairs = 4
    elif metadataRows > 2 * maxMetadataRowsPerColumn:
        metadataColumnPairs = 3
    elif metadataRows > maxMetadataRowsPerColumn:
        metadataColumnPairs = 2
    else:
        metadataColumnPairs = 1

    columns = metadataColumnPairs * 2
    rows = min(maxMetadataRowsPerColumn, metadataRows)

    # Get the rectangle the content will draw in
    contentLeft, contentWidth, contentTop, contentHeight = getContentRect(
        presentation, slide, slideTitleBottom, tableMargin
    )

    tableHeight = min(contentHeight, Inches(0.25) * rows)

    # Figure out the width of a single-width column
    columnWidthUnit = int(contentWidth / (2 * metadataColumnPairs))

    # Create the table with the above number of rows and columns
    newTable = slide.shapes.add_table(
        rows, columns, tableMargin, contentTop, contentWidth, tableHeight
    ).table

    # Don't want headings
    newTable.first_row = False

    cols = newTable.columns
    for cp in range(metadataColumnPairs):
        cols[2 * cp].width = columnWidthUnit
        cols[2 * cp + 1].width = columnWidthUnit

    row = 0
    column = 0
    for item in metadata:
        key, value = item

        if row == maxMetadataRowsPerColumn:
            # Move to next column
            column += 2
            row = 0

        # Set text of metadata key cell
        newTable.cell(row, column).text = key
        if md2pptx.globals.processingOptions.dynamicallyChangedOptions.get(key) is not None:
            # Set text of metadata value cell - with asterisk
            newTable.cell(row, column + 1).text = value + "*"

            # Colour key cell blue
            p1 = newTable.cell(row, column).text_frame.paragraphs[0]
            p1.font.color.rgb = RGBColor.from_string("0000FF")

            # Colour value cell blue
            p2 = newTable.cell(row, column + 1).text_frame.paragraphs[0]
            p2.font.color.rgb = RGBColor.from_string("0000FF")
        else:
            # Set text of metadata value cell - without asterisk
            newTable.cell(row, column + 1).text = value

        if metadataColumnPairs == 5:
            newTable.cell(row, column).text_frame.paragraphs[0].font.size = Pt(8)
            newTable.cell(row, column + 1).text_frame.paragraphs[0].font.size = Pt(8)
        elif metadataColumnPairs == 4:
            newTable.cell(row, column).text_frame.paragraphs[0].font.size = Pt(10)
            newTable.cell(row, column + 1).text_frame.paragraphs[0].font.size = Pt(10)
        elif metadataColumnPairs == 3:
            newTable.cell(row, column).text_frame.paragraphs[0].font.size = Pt(12)
            newTable.cell(row, column + 1).text_frame.paragraphs[0].font.size = Pt(12)
        elif metadataColumnPairs == 2:
            newTable.cell(row, column).text_frame.paragraphs[0].font.size = Pt(14)
            newTable.cell(row, column + 1).text_frame.paragraphs[0].font.size = Pt(14)
        else:
            newTable.cell(row, column).text_frame.paragraphs[0].font.size = Pt(16)
            newTable.cell(row, column + 1).text_frame.paragraphs[0].font.size = Pt(16)
        row += 1


# Note: This doesn't use formatTitle()
def createTitleOrSectionSlide(
    presentation,
    slideNumber,
    titleText,
    layout,
    titleSize,
    subtitleText,
    subtitleSize,
    notes_text,
):
    marginBase = md2pptx.globals.processingOptions.getCurrentOption("marginBase")

    slide = addSlide(presentation, presentation.slide_layouts[layout], None)

    # Add title
    title = findTitleShape(slide)
    flattenedTitle = addFormattedText(title.text_frame.paragraphs[0], titleText)

    title.text_frame.paragraphs[0].font.size = Pt(titleSize)

    # Add subtitle - if there is one
    if (subtitleText.strip() != "") & (subtitleText[0:2] != "\n\n"):
        # There is a subtitle
        chunks = subtitleText.strip().split("\n\n")
        subtitleText = chunks[0]
        notes_text = "\n\n".join(chunks[1:])
        createSlideNotes(slide, notes_text)

        subtitleShape = findBodyShape(slide)
        if subtitleShape != None:
            addFormattedText(subtitleShape.text_frame.paragraphs[0], subtitleText)
            subtitleShape.text_frame.paragraphs[0].font.size = Pt(subtitleSize)
        else:
            print("Warning: No subtitle shape on this slide to add text to.")
    else:
        # Reformat subtitle shape to be out of the way
        subtitleShape = findBodyShape(slide)
        if subtitleShape != None:
            subtitleShape.top = title.top + title.height + marginBase * 2
            subtitleShape.width = title.width
            subtitleShape.left = title.left
            subtitleShape.height = marginBase * 2

    reportSlideTitle(slideNumber, 1, flattenedTitle)

    if want_numbers_headings is True:
        addFooter(presentation, slideNumber, slide)

    return slide


def handleWhateverGraphicType(GraphicFilename):
    # Handles both physical file and URI file types
    GraphicFilename = GraphicFilename.strip()

    # First ensure we have the data in a file and know if the source was a URI
    if ":" in GraphicFilename:
        # Is a URI - so we have to retrieve it and store it in a temporary file
        is_uri = True

        # Massage the URI into a printable filename
        if len(GraphicFilename) > 50:
            printableGraphicFilename = (
                GraphicFilename[:25] + "..." + GraphicFilename[-25:]
            )
        else:
            printableGraphicFilename = GraphicFilename

        # Get Temporary File Directory - which might be None
        tempDir = md2pptx.globals.processingOptions.getCurrentOption("tempDir")

        # Retrieve the data into a temporary file
        try:
            operUrl = urllib.request.urlopen(GraphicFilename)

        except urllib.error.HTTPError as e:
            print("HTTP error: " + str(e.code))
            return GraphicFilename, printableGraphicFilename

        except socket.error as s:
            print("Socket error. (Web site not found).")
            return GraphicFilename, printableGraphicFilename

        data = operUrl.read()

        # Get Content-Type header - if set
        content_type = ""

        if str(type(operUrl)) == "<class 'http.client.HTTPResponse'>":
            # Can try to get content header
            content_type = operUrl.getheader("content-type")

            if content_type == "":
                # Obtain file extension by searching in the URL
                extensionPos = GraphicFilename.rindex(".")
                lastSlashPos = GraphicFilename.rindex("/")
                if lastSlashPos > extensionPos:
                    fileExt = ""

                else:
                    fileExt = GraphicFilename[extensionPos:]

            else:
                # Set file extension based on Content-Type header_name
                # Note: It's been translated to lower case above
                if content_type in ["image/jpeg", "image/jpg"]:
                    # Note: only the first of these two is legitimate
                    fileExt = "jpg"
                elif content_type == "image/png":
                    fileExt = "png"
                elif content_type in ["image/svg+xml", "image/svg"]:
                    # Note: only the first of these two is legitimate
                    fileExt = "svg"
                elif content_type == "application/postscript":
                    fileExt = "eps"
                else:
                    fileExt = None
        else:
            fileExt = None

        # Store in a temporary file
        try:
            tempGraphicFile = tempfile.NamedTemporaryFile(
                delete=False, suffix=fileExt, dir=tempDir
            )
        except IOError as e:
            print("Couldn't create temporary file. md2pptx terminating")
            exit()

        tempGraphicFile.write(data)
        convertibleFilename = tempGraphicFile.name
        tempGraphicFile.close()

    else:
        is_uri = False

        # Files don't get their names edited
        printableGraphicFilename = GraphicFilename
        convertibleFilename = GraphicFilename

    if is_uri:
        lastSlash = GraphicFilename.rfind("/")
        lastDot = GraphicFilename.rfind(".")

        PNGname = GraphicFilename[lastSlash + 1 : lastDot] + ".PNG"
    else:
        PNGname = GraphicFilename[:-4] + ".PNG"

    # Process the file - whatever the origin - based on file extension
    if ".svg" in GraphicFilename.lower():
        # is an SVG file
        if have_cairosvg:
            # Convert SVG file to temporary PNG
            # Store in a temporary file

            # Get Temporary File Directory - which might be None
            tempDir = md2pptx.globals.processingOptions.getCurrentOption("tempDir")

            try:
                graphicFile = tempfile.NamedTemporaryFile(
                    delete=False, suffix=".PNG", dir=tempDir
                )
            except IOError as e:
                print("Couldn't create temporary file. md2pptx terminating")
                exit()

            cairosvg.svg2png(file_obj=open(convertibleFilename), write_to=graphicFile)

            # Retrieve the temporary file name
            GraphicFilename = graphicFile.name

            if md2pptx.globals.processingOptions.getCurrentOption("exportGraphics"):
                try:
                    shutil.copy(GraphicFilename, PNGname)
                except:
                    print("Copy error: " + PNGname)

        else:
            print("Don't have CairoSVG installed. Terminating.")
            sys.exit()
    elif ".eps" in GraphicFilename.lower():
        if have_pillow:
            # Get EPS file
            im = PIL.Image.open(GraphicFilename)

            # Get Temporary File Directory - which might be None
            tempDir = md2pptx.globals.processingOptions.getCurrentOption("tempDir")

            # Store in a temporary file
            try:
                graphicFile = tempfile.NamedTemporaryFile(
                    delete=False, suffix=".PNG", dir=tempDir
                )
            except IOError as e:
                print("Couldn't create temporary file. md2pptx terminating")
                exit()

            try:
                im.save(graphicFile)
            except:
                print("Could not convert EPS file. Is Ghostview installed?\n")
                print("Terminating.\n")
                sys.exit()

            # Retrieve the temporary file name
            GraphicFilename = graphicFile.name
            if md2pptx.globals.processingOptions.getCurrentOption("exportGraphics"):
                try:
                    shutil.copy(GraphicFilename, PNGname)
                except:
                    print("Copy error: " + PNGname)

    else:
        GraphicFilename = convertibleFilename

    return GraphicFilename, printableGraphicFilename


def handleGraphViz(slide, renderingRectangle, codeLines, codeType):
    # Condition GraphViz source
    s = graphviz.Source("\n".join(codeLines), format="png")

    # Invent a temporary filename for the rendered graphic
    dotFile = "md2pptx-temporary-dotfile.png"

    # Render the .dot source as a graphic
    s.render(cleanup=True, outfile=dotFile)

    # Figure out the dimensions of the rendered graphic
    dotGraphicWidth, dotGraphicHeight = getGraphicDimensions(dotFile)

    # Adjust those dimensions with the usual scaling rules
    (dotPicWidth, dotPicHeight, scaledByHeight) = scalePicture(
        renderingRectangle.width,
        renderingRectangle.height,
        dotGraphicWidth,
        dotGraphicHeight,
    )

    # Add the picture to the current slide
    slide.shapes.add_picture(
        dotFile,
        renderingRectangle.left + (renderingRectangle.width - dotPicWidth) / 2,
        renderingRectangle.top + (renderingRectangle.height - dotPicHeight) / 2,
        dotPicWidth,
        dotPicHeight,
    )

    # Delete the temporary graphic file
    os.remove(dotFile)


def handleFunnel(slide, renderingRectangle, codeLines, codeType):
    funnelColours = md2pptx.globals.processingOptions.getCurrentOption("funnelColours")
    funnelBorderColour = md2pptx.globals.processingOptions.getCurrentOption("funnelBorderColour")
    funnelTitleColour = md2pptx.globals.processingOptions.getCurrentOption("funnelTitleColour")
    funnelTextColour = md2pptx.globals.processingOptions.getCurrentOption("funnelTextColour")
    funnelLabelsPercent = md2pptx.globals.processingOptions.getCurrentOption("funnelLabelsPercent")
    funnelLabelPosition = md2pptx.globals.processingOptions.getCurrentOption("funnelLabelPosition")
    funnelWidest = md2pptx.globals.processingOptions.getCurrentOption("funnelWidest")

    f = md2pptx.funnel.Funnel()

    f.makeFunnel(
        slide,
        renderingRectangle,
        codeLines,
        funnelColours,
        codeType,
        funnelBorderColour,
        funnelTitleColour,
        funnelTextColour,
        funnelLabelsPercent,
        funnelLabelPosition,
        funnelWidest,
    )

# Handler function for immediately executing python in a code block
def handleRunPython(pythonType, prs, slide, renderingRectangle, codeLinesOrFile, codeType):
    r = md2pptx.runPython.RunPython()

    if pythonType == "inline":
        r.run(prs, slide, renderingRectangle, codeLinesOrFile, codeType)
    else:
        r.runFromFile(codeLinesOrFile[0], prs, slide, renderingRectangle)

def createCodeBlock(slideInfo, slide, renderingRectangle, codeBlockNumber):
    monoFont = md2pptx.globals.processingOptions.getCurrentOption("monoFont")
    baseTextSize = md2pptx.globals.processingOptions.getCurrentOption("baseTextSize")
    defaultBaseTextSize = md2pptx.globals.processingOptions.getDefaultOption("baseTextSize")

    # A variable number of newlines appear before the actual code
    codeLines = slideInfo.code[codeBlockNumber]

    # Figure out code slide type
    if codeLines[0].startswith("<pre"):
        codeType = "pre"
    elif codeLines[0].startswith("<code"):
        codeType = "code"
    elif codeLines[0].startswith("```"):
        if codeLines[0] == "```":
            # Plain backticks
            codeType = "backticks"
        else:
            # Some other type - such as GraphViz .dot
            codeTypeRaw = codeLines[0][3:].strip()
            codeType = codeTypeRaw.lower()

            if codeType.startswith("dot"):
                # GraphViz dot
                if have_graphviz is False:
                    # Don't have GraphViz so warn and treat as basic backticks
                    sys.stderr.write(
                        "GraphViz not installed. Rendering as backticks.\n"
                    )

                    codeType = "backticks"
                else:
                    # Have GraphViz installed so use it ti add a picture
                    handleGraphViz(slide, renderingRectangle, codeLines[1:-1], codeType)

                    return slide

            elif codeType.startswith("funnel"):
                # Built-in Funnel Diagram
                handleFunnel(slide, renderingRectangle, codeLines[1:-1], codeType)

                return slide

            elif codeType.startswith("run-python"):
                if codeType == "run-python":
                    # Run inline lines as Python
                    handleRunPython("inline", prs, slide, renderingRectangle, codeLines[1:-1], codeType)
                else:
                    # Run Python file
                    runPythonFile = codeTypeRaw.split(" ")[1]
                    handleRunPython("fromfile", prs, slide, renderingRectangle, [runPythonFile], codeType)

                return slide

            else:
                # Some other type with backticks
                codeType = "backticks"

    else:
        codeType = "indented"

    # Handle any trailing empty lines
    while codeLines[-1] == "":
        codeLines.pop(-1)

    # Handle any leading <pre>, <code>, triple backtick line
    if startswithOneOf(codeLines[0], ["<pre>", "<code>", "```"]):
        codeLines.pop(0)

    # Handle any trailing </pre>, </code>, triple backtick line
    if startswithOneOf(codeLines[-1], ["</pre>", "</code>", "```"]):
        codeLines.pop(-1)

    codeBox = slide.shapes.add_textbox(
        renderingRectangle.left,
        renderingRectangle.top,
        renderingRectangle.width,
        renderingRectangle.height,
    )

    # Try to control text frame but SHAPE_TO_FIT_TEXT doesn't seem to work
    tf = codeBox.text_frame
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    tf.word_wrap = False

    # Fill the code box with background colour - whether explicit or defaulted
    fill = codeBox.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(
        md2pptx.globals.processingOptions.getCurrentOption("codeBackground")
    )

    # Get the sole paragraph
    p = codeBox.text_frame.paragraphs[0]

    # Set the font size slightly smaller than usual
    if len(codeLines) >= 20:
        divisor = 1.5
    else:
        divisor = 1.2
    if baseTextSize > 0:
        p.font.size = int(Pt(baseTextSize) / divisor)
    else:
        p.font.size = int(Pt(defaultBaseTextSize) / divisor)

    # Estimate how wide the code box would need to be at the current font size
    # versus actual width
    codeColumns = md2pptx.globals.processingOptions.getCurrentOption("codeColumns")
    fixedPitchHeightWidthRatio = md2pptx.globals.processingOptions.getCurrentOption(
        "fixedPitchHeightWidthRatio"
    )

    estimatedWidthVersusCodeboxWidth = (
        p.font.size * codeColumns / codeBox.width / fixedPitchHeightWidthRatio
    )
    if estimatedWidthVersusCodeboxWidth > 1:
        # The code is wider so shrink the font so the code fits
        p.font.size = p.font.size / estimatedWidthVersusCodeboxWidth
    else:
        # The code is narrower so shrink the code textbox so the code just fits
        # - assuming declared width is accurate
        codeBox.width = int(p.font.size * codeColumns / fixedPitchHeightWidthRatio)

        # Center the code box - actually don't - 5 October 2021 temporary "fix"
        # codeBox.left = int((presentation.slide_width - codeBox.width) / 2)

    # Use the code foreground colour - whether explicit or defaulted
    p.font.color.rgb = RGBColor.from_string(
        md2pptx.globals.processingOptions.getCurrentOption("codeforeground")
    )

    # Adjust code box height based on lines
    codeBox.height = min(
        len(codeLines) * Pt(baseTextSize + 5), renderingRectangle.height
    )

    # Add code
    if codeType == "pre":
        # span elements want special handling
        for codeLine in codeLines:
            # Resolve eg entity references
            codeLine = resolveSymbols(codeLine)

            # Split the line - and maybe there are spans
            spanFragments = codeLine.split("<span ")
            if len(spanFragments) > 1:
                textArray = []
                # Break line down into what will become runs
                for fragmentNumber, fragment in enumerate(spanFragments):
                    if fragmentNumber > 0:
                        # Find start of span class
                        fragment = "<span " + fragment
                        if spanClassMatch := md2pptx.globals.spanClassRegex.match(fragment):
                            afterSpanTag = fragment[spanClassMatch.end(1) :]
                            className = afterSpanTag[7 : afterSpanTag.index(">") - 1]
                            if (
                                (className in md2pptx.globals.bgcolors)
                                | (className in md2pptx.globals.fgcolors)
                                | (className in md2pptx.globals.emphases)
                            ):
                                afterClosingAngle = afterSpanTag[
                                    afterSpanTag.index(">") + 1 :
                                ]
                                startEnd = afterClosingAngle.index("</span>")
                                afterSpan2 = afterClosingAngle[:startEnd]
                                afterSpan3 = afterClosingAngle[startEnd + 7 :]
                                textArray.append(["SpanClass", [className, afterSpan2]])
                                textArray.append(["Normal", afterSpan3])
                                fragment = ""
                            else:
                                print(
                                    className
                                    + " is not defined. Ignoring reference to it in <span> element."
                                )
                        elif spanStyleMatch := md2pptx.globals.spanStyleRegex.match(fragment):
                            afterSpanTag = fragment[spanStyleMatch.end(1) :]
                            styleText = afterSpanTag[7 : afterSpanTag.index(">") - 1]
                            styleElements = styleText.split(";")
                            afterClosingAngle = afterSpanTag[
                                afterSpanTag.index(">") + 1 :
                            ]
                            startEnd = afterClosingAngle.index("</span>")
                            afterSpan2 = afterClosingAngle[:startEnd]
                            afterSpan3 = afterClosingAngle[startEnd + 7 :]
                            textArray.append(["SpanStyle", [styleText, afterSpan2]])
                            textArray.append(["Normal", afterSpan3])
                            fragment = ""
                    else:
                        textArray.append(["Normal", fragment])

                # Now we have a text array we can add the runs for the line
                for textArrayItem in textArray:
                    textArrayItemType = textArrayItem[0]
                    if textArrayItemType == "Normal":
                        # Is not in a span element bracket
                        className = ""
                        spanStyle = ""
                        spanText = textArrayItem[1]

                    elif textArrayItemType == "SpanClass":
                        # Is in a span class element bracket
                        className = textArrayItem[1][0]
                        spanText = textArrayItem[1][1]
                        spanStyle = ""

                    else:
                        # Is in a span style element bracket
                        spanStyle = textArrayItem[1][0]
                        spanText = textArrayItem[1][1]

                    if spanText != "":
                        run = p.add_run()
                        run.text = spanText
                        font = run.font
                        font.name = monoFont

                    if className != "":
                        # Augment run with whatever the span class calls for
                        handleSpanClass(run, className)

                    if spanStyle != "":
                        # Augment the run with whatever the style calls for
                        handleSpanStyle(run, spanStyle)

                # Add terminating newline
                run = p.add_run()
                run.text = "\n"
                font = run.font
                font.name = monoFont
            else:
                # Line has no spans in
                run = p.add_run()
                run.text = codeLine + "\n"
                font = run.font
                font.name = monoFont

    else:
        # span doesn't need treating specially
        for codeLine in codeLines:
            # Resolve eg entity references
            codeLine = resolveSymbols(codeLine)

            run = p.add_run()
            run.text = codeLine + "\n"
            font = run.font
            font.name = monoFont

    return slide


def createAbstractSlide(presentation, slideNumber, titleText, paragraphs):
    titleOnlyLayout = md2pptx.globals.processingOptions.getCurrentOption("titleOnlyLayout")
    marginBase = md2pptx.globals.processingOptions.getCurrentOption("marginBase")
    pageTitleSize = md2pptx.globals.processingOptions.getCurrentOption("pageTitleSize")
    pageSubtitleSize = md2pptx.globals.processingOptions.getCurrentOption("pageSubtitleSize")

    slide = addSlide(presentation, presentation.slide_layouts[titleOnlyLayout], None)

    shapes = slide.shapes

    # Add title and constrain its size and placement
    slideTitleBottom, title, flattenedTitle = formatTitle(
        presentation, slide, titleText, pageTitleSize, pageSubtitleSize
    )

    reportSlideTitle(slideNumber, 3, "Abstract: " + flattenedTitle)

    # Get the rectangle the content will draw in
    contentLeft, contentWidth, contentTop, contentHeight = getContentRect(
        presentation, slide, slideTitleBottom, marginBase
    )

    # Add abstract text
    abstractBox = slide.shapes.add_textbox(
        contentLeft,
        contentTop,
        contentWidth,
        contentHeight,
    )

    p = abstractBox.text_frame.paragraphs[0]
    tf = abstractBox.text_frame
    f = p.font
    f.size = Pt(22)
    for para, abstractParagraph in enumerate(paragraphs):
        paragraphLevel, paragraphText, paragraphType = abstractParagraph

        if para > 0:
            # Spacer paragraph
            p = tf.add_paragraph()
            f = p.font
            f.size = Pt(22)

            # Content paragraph
            p = tf.add_paragraph()
            f = p.font
            f.size = Pt(22)
        addFormattedText(p, paragraphText)

    tf.word_wrap = True

    if want_numbers_content is True:
        addFooter(presentation, slideNumber, slide)

    return slide


# Unified creation of a table or a code or a content slide
def createContentSlide(presentation, slideNumber, slideInfo):
    titleOnlyLayout = md2pptx.globals.processingOptions.getCurrentOption("titleOnlyLayout")
    contentSlideLayout = md2pptx.globals.processingOptions.getCurrentOption("contentSlideLayout")
    marginBase = md2pptx.globals.processingOptions.getCurrentOption("marginBase")
    pageTitleSize = md2pptx.globals.processingOptions.getCurrentOption("pageTitleSize")
    pageSubtitleSize = md2pptx.globals.processingOptions.getCurrentOption("pageSubtitleSize")

    # slideInfo's body text is only filled in if there is code - and that's
    # where the code - plus preamble and postamble is.
    if slideInfo.code != "":
        haveCode = True
    else:
        haveCode = False

    # Create the slide and check for bullets and/or cards
    if (slideInfo.bullets == []) & (slideInfo.cards == []):
        # No bullets or cards so "title only"
        slideLayout = titleOnlyLayout
        haveBulletsCards = False
    else:
        # Either bullets or cards or both so not "title only"
        slideLayout = contentSlideLayout
        haveBulletsCards = True

    slide = addSlide(presentation, presentation.slide_layouts[slideLayout], slideInfo)

    # Check for table / graphics content
    if slideInfo.tableRows == []:
        haveTableGraphics = False
    else:
        haveTableGraphics = True

    ####################################################################
    # At this point haveCode, haveBulletsCards, haveTableGraphics have #
    # been appropriately set                                           #
    ####################################################################

    # Add slide title
    titleText = slideInfo.titleText

    slideTitleBottom, title, flattenedTitle = formatTitle(
        presentation, slide, titleText, pageTitleSize, pageSubtitleSize
    )

    # Log slide's title
    reportSlideTitle(slideNumber, 3, flattenedTitle)

    ####################################################################
    # Get the dimensions of the content area to place all content in   #
    ####################################################################
    contentLeft, contentWidth, contentTop, contentHeight = getContentRect(
        presentation, slide, slideTitleBottom, marginBase
    )

    ####################################################################
    # Check whether there are too many elements in the sequence to     #
    # render - and warn if there are. Then calculate how many to render#
    ####################################################################
    if len(slideInfo.sequence) > maxBlocks:
        print(f"Too many blocks to render. Only {str(maxBlocks)} will be rendered.")
    blocksToRender = min(maxBlocks, len(slideInfo.sequence))

    ####################################################################
    # Get the dimensions of the rectangles we'll place the graphics in #
    # and their top left corner coordinates                            #
    ####################################################################
    allContentSplit = 0
    contentSplit = md2pptx.globals.processingOptions.getCurrentOption("contentSplit")
    for b in range(blocksToRender):
        allContentSplit = allContentSplit + contentSplit[b]

    verticalCursor = contentTop
    horizontalCursor = contentLeft

    codeBlockNumber = 0
    tableBlockNumber = 0

    for b in range(blocksToRender):
        if md2pptx.globals.processingOptions.getCurrentOption("contentSplitDirection") == "vertical":
            # Height and top
            blockHeight = int(contentHeight * contentSplit[b] / allContentSplit)
            blockTop = verticalCursor
            verticalCursor = verticalCursor + blockHeight

            # Width and left
            blockWidth = contentWidth
            blockLeft = contentLeft
        else:
            # Height and top
            blockHeight = contentHeight
            blockTop = contentTop

            # Width and left
            blockWidth = int(contentWidth * contentSplit[b] / allContentSplit)
            blockLeft = horizontalCursor
            horizontalCursor = horizontalCursor + blockWidth

        renderingRectangle = Rectangle(blockTop, blockLeft, blockHeight, blockWidth)

        if slideInfo.sequence[b] == "table":
            createTableBlock(slideInfo, slide, renderingRectangle, tableBlockNumber)
            tableBlockNumber += 1

        elif slideInfo.sequence[b] == "list":
            createListBlock(slideInfo, slide, renderingRectangle)
        else:
            createCodeBlock(slideInfo, slide, renderingRectangle, codeBlockNumber)
            codeBlockNumber += 1

    if want_numbers_content is True:
        addFooter(presentation, slideNumber, slide)

    return slide


def createListBlock(slideInfo, slide, renderingRectangle):
    horizontalCardGap = md2pptx.globals.processingOptions.getCurrentOption("horizontalcardgap")
    verticalCardGap = md2pptx.globals.processingOptions.getCurrentOption("verticalcardgap")
    cardTitleAlign = md2pptx.globals.processingOptions.getCurrentOption("cardtitlealign")
    cardTitlePosition = md2pptx.globals.processingOptions.getCurrentOption("cardtitleposition")
    cardShape = md2pptx.globals.processingOptions.getCurrentOption("cardshape")
    cardLayout = md2pptx.globals.processingOptions.getCurrentOption("cardlayout")
    cardPercent = md2pptx.globals.processingOptions.getCurrentOption("cardpercent")
    cardShadow = md2pptx.globals.processingOptions.getCurrentOption("cardshadow")
    cardTitleSize = md2pptx.globals.processingOptions.getCurrentOption("cardtitlesize")
    cardBorderWidth = md2pptx.globals.processingOptions.getCurrentOption("cardborderwidth")
    cardBorderColour = md2pptx.globals.processingOptions.getCurrentOption("cardbordercolour")
    cardTitleColour = md2pptx.globals.processingOptions.getCurrentOption("cardtitlecolour")
    cardTitleBackgrounds = md2pptx.globals.processingOptions.getCurrentOption("cardtitlebackground")
    cardColours = md2pptx.globals.processingOptions.getCurrentOption("cardcolour")
    cardDividerColour = md2pptx.globals.processingOptions.getCurrentOption("carddividercolour")
    cardGraphicSize = md2pptx.globals.processingOptions.getCurrentOption("cardgraphicsize")
    cardGraphicPosition = md2pptx.globals.processingOptions.getCurrentOption("cardGraphicPosition")
    cardGraphicPadding = int(Inches(md2pptx.globals.processingOptions.getCurrentOption("cardgraphicpadding")))
    marginBase = md2pptx.globals.processingOptions.getCurrentOption("marginBase")
    pageTitleSize = md2pptx.globals.processingOptions.getCurrentOption("pageTitleSize")
    pageSubtitleSize = md2pptx.globals.processingOptions.getCurrentOption("pageSubtitleSize")

    # Get bulleted text shape - either for bullets above cards or first card's body shape
    bulletsShape = findBodyShape(slide)

    # Set bulleted shape top, left, width
    bulletsShape.top = renderingRectangle.top
    bulletsShape.left = renderingRectangle.left
    bulletsShape.width = renderingRectangle.width

    bulletCount = len(slideInfo.bullets)

    # Set bulleted text height - depending on whether there's a card
    # Remainder is card area height - if there are cards
    if slideInfo.cards == []:
        # There are no cards so the bullets shape takes the whole content area
        bulletsShape.height = renderingRectangle.height

        # There are no cards so the card area is zero height
        cardAreaHeight = 0
        cardCount = 0
    else:
        # There are cards
        if bulletCount > 0:
            # Bullets shape vertically shortened
            bulletsShape.height = int(
                renderingRectangle.height * (100 - cardPercent) / 100
            )

            # Card area takes the rest of the content area
            cardAreaHeight = int(renderingRectangle.height) - bulletsShape.height
        else:
            # No bullets so content is all cards
            bulletsShape.height = 0

            cardAreaHeight = renderingRectangle.height

        cardCount = len(slideInfo.cards)

        ###########################################################
        # Work out card dimensions - based on the various layouts #
        ###########################################################

        # card width applies to card title, card graphic, card background, card body
        if cardLayout == "horizontal":
            # Divide horizontal card space up
            cardWidth = int(
                (renderingRectangle.width - Inches(horizontalCardGap) * (cardCount - 1))
                / cardCount
            )
        else:
            # Card takes all the horizontal space
            cardWidth = int(renderingRectangle.width)

        # Calculate title top and height - horizontal layout
        if cardTitleSize > 0:
            # Specified by user. "72" because in points
            cardTitleHeightRaw = Inches(cardTitleSize / 72)
        else:
            # Shrunk to 2/3 of page title height.  "72" because in points
            cardTitleHeightRaw = Inches(int(10000 * pageTitleSize * 2 / 3 / 72) / 10000)

        # Adjust title height to be slightly larger than the text
        cardTitleHeight = cardTitleHeightRaw + Inches(0.1)

        cardGraphicSizeRaw = int(Inches(cardGraphicSize))

        if bulletCount > 0:
            # Bullets so cards and their titles occupy less than whole height
            cardAreaTop = bulletsShape.height + renderingRectangle.top

        else:
            # No bullets so cards and their titles occupy whole height
            cardAreaTop = renderingRectangle.top

        if cardLayout == "horizontal":
            # Card takes up all the card area, vertically
            cardHeight = cardAreaHeight
        else:
            # Card layout is horizontol so card height is just a proportion
            # of the card area height

            if cardTitlePosition == "above":
                paddingFactor = Inches(verticalCardGap - 0.05)
            else:
                paddingFactor = Inches(verticalCardGap)

            cardHeight = int((cardAreaHeight) / cardCount - paddingFactor)

        # Store slide title shape for cloning
        slideTitleShape = findTitleShape(slide)

        ###############################################################
        # Work out whether any card has a printable title. If not set #
        # cardTitleHeight to 0                                        #
        ###############################################################
        cardWithPrintableTitle = False

        for c, card in enumerate(slideInfo.cards):
            # Check if there are any titles for any of the cards
            if card.title != "&nbsp;":
                cardWithPrintableTitle = True

        if not cardWithPrintableTitle:
            # Zero card title height
            cardTitleHeight = 0

        ###########################################################
        # Work out card positions - based on the various layouts  #
        ###########################################################
        for c, card in enumerate(slideInfo.cards):
            # Calculate each card's vertical position
            if cardLayout == "horizontal":
                # Card top is at top of card area
                card.top = cardAreaTop
            else:
                # vertical so card tops are progressively further down the card
                # area
                card.top = int((cardHeight + paddingFactor) * c + cardAreaTop)

            # Calculate each card's background and body top
            if cardTitlePosition == "above":
                # Card title (if any) above card background - so background top is
                # below card top
                card.backgroundTop = card.top + cardTitleHeight
                card.bodyTop = card.backgroundTop
            else:
                # card title (if any) inside card background - so background top is
                # card top
                card.backgroundTop = card.top

                # Leave room above the card body for the card title (if any)
                card.bodyTop = card.backgroundTop + cardTitleHeight

            # Calculate each card's horizontal position
            if cardLayout == "horizontal":
                # Card lefts are progressively across the card area
                card.left = marginBase + c * (cardWidth + Inches(horizontalCardGap))
            else:
                # Vertical so card lefts are at the left of the card area
                card.left = marginBase

            # Card title modeled on slide title - but smaller
            cardTitleShape = addClonedShape(slide, slideTitleShape)

            card.titleShape = cardTitleShape

            if card.graphic != "":
                # This card has a graphic so get its name
                card.graphicDimensions = getGraphicDimensions(card.graphic)

                # Create card graphic shape - to be resized later
                card.graphicShape = slide.shapes.add_picture(
                    card.graphic,
                    Inches(0),
                    Inches(0),
                )

                reportGraphicFilenames(card.printableFilename)

                if (card.mediaURL is not None) | (card.graphicTitle is not None):
                    mediaURL = card.mediaURL

                    graphicTitle = "" if card.graphicTitle == None else card.graphicTitle

                    pictureInfos.append(
                        (card.graphicShape, mediaURL, graphicTitle)
                    )

            elif card.mediaInfo is not None:
                # This card has a video so get its dimensions etc
                card.mediaDimensions = getVideoInfo(card.mediaInfo)

                # Create card video shape - to be resized later
                card.mediaShape = slide.shapes.add_movie(
                    card.mediaInfo.source,
                    Inches(0),
                    Inches(0),
                    Inches(0),
                    Inches(0),
                    card.mediaInfo.poster
                )

                reportGraphicFilenames(card.printableFilename)

            else:
                # Some of this is probably not needed
                card.graphicDimensions = None
                card.mediaDimensions = None

            # Clear text from cloned title and add in the title text
            cardTitleShape.text_frame.paragraphs[0].text = ""
            addFormattedText(cardTitleShape.text_frame.paragraphs[0], card.title)

            # Set card title font size
            if cardTitleSize > 0:
                # Explicitly set title size
                cardTitleShape.text_frame.paragraphs[0].font.size = Pt(cardTitleSize)
            else:
                # Not explicitly set - so default to 2/3 slide title size
                cardTitleShape.text_frame.paragraphs[0].font.size = Pt(
                    pageTitleSize * 2 / 3
                )

            # Titles are aligned one of three ways
            if cardTitleAlign == "l":
                cardTitleShape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            elif cardTitleAlign == "c":
                cardTitleShape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            else:
                cardTitleShape.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

            # Fill in card's background - if necessary
            if (cardTitleBackgrounds[0] != ("None", "")) & (
                cardTitlePosition != "inside"
            ):
                # Card title background picked round-robin from array
                cardTitleBackground = cardTitleBackgrounds[
                    c % len(cardTitleBackgrounds)
                ]

                fill = cardTitleShape.fill
                fill.solid()

                setColour(fill.fore_color, cardTitleBackground)

            # Create card background and make sure it's behind the card body
            # (and maybe card title)
            if cardShape == "rounded":
                # Rounded Rectangle for card
                cardBackgroundShape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(0),
                    Inches(0),
                    Inches(0),
                    Inches(0),
                )

                # Rounding adjustment works better with different values for horizontal and vertical cards
                if cardLayout == "horizontal":
                    # Make the rounding radius small. This is 1/4 the default
                    cardBackgroundShape.adjustments[0] = 0.0416675
                else:
                    # Make the rounding radius smallish. This is 1/2 the default
                    cardBackgroundShape.adjustments[0] = 0.083335
            else:
                # Squared-corner Rectangle for card
                cardBackgroundShape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0), Inches(0)
                )

            card.backgroundShape = cardBackgroundShape

            if cardShape == "line":
                # Ensure no fill for card background
                cardBackgroundShape.fill.background()

            # Ensure card background is at the back
            sendToBack(slide.shapes, cardBackgroundShape)

            # Card shape modeled on bulleted list
            if (bulletCount > 0) | (c > 0):
                # Copy the bullets shape for not-first card body shape
                cardBodyShape = addClonedShape(slide, bulletsShape)
            else:
                # Co-opt bullets shape as first card shape
                cardBodyShape = bulletsShape

            card.bodyShape = cardBodyShape

            # Make card's body transparent
            fill = cardBodyShape.fill
            fill.background()

            # Fill in card's background - if necessary
            if (cardColours[0] != ("None", "")) & (cardShape != "line"):
                # Card background volour picked round-robin from array
                cardColour = cardColours[c % len(cardColours)]

                fill = cardBackgroundShape.fill
                fill.solid()

                setColour(fill.fore_color, cardColour)

    #######################################################################
    # Adjust bullets shape height - and calculate verticals for any cards #
    #######################################################################

    # Set bottom of bullets shape
    bulletsShape.bottom = bulletsShape.top + bulletsShape.height

    # Fill in the main bullets shape
    renderText(bulletsShape, slideInfo.bullets)

    # Second go on each card
    for c, card in enumerate(slideInfo.cards):
        # Get the shapes for this card - including any graphic
        cardBackgroundShape = card.backgroundShape
        cardTitleShape = card.titleShape
        cardBodyShape = card.bodyShape
        cardGraphicShape = card.graphicShape
        cardMediaShape = card.mediaShape

        # Get dimensions of any graphic
        if cardGraphicShape is not None :
            cardMediaNativeWidth, cardMediaNativeHeight = card.graphicDimensions

        elif cardMediaShape is not None:
            cardMediaNativeWidth, cardMediaNativeHeight, _, _ = card.mediaDimensions

        else:
            cardMediaNativeWidth, cardMediaNativeHeight = (0, 0)

        # Set the card shapes' width
        cardBackgroundShape.width = cardWidth
        if cardLayout == "horizontal":
            cardBodyShape.width = cardWidth
        else:
            cardBodyShape.width = cardWidth - cardGraphicSizeRaw

        cardTitleShape.width = cardWidth

        # Set the card shapes' left side
        cardBackgroundShape.left = card.left
        cardTitleShape.left = card.left
        if (cardLayout == "horizontal") | (cardGraphicPosition == "after"):
            cardBodyShape.left = card.left
        else:
            cardBodyShape.left = card.left + cardGraphicSizeRaw + \
                2 * cardGraphicPadding

        # Position card title
        cardTitleShape.top = card.top
        cardTitleShape.height = cardTitleHeight

        # Colour the title - if cardTitleColour specified
        if cardTitleColour != ("None", ""):
            setColour(
                cardTitleShape.text_frame.paragraphs[0].font.color, cardTitleColour
            )

        # Calculate positions and heights within card background of body
        if cardTitlePosition == "above":
            # Any card titles would be above the rest of the card
            cardBackgroundShape.top = card.top + cardTitleHeight
            cardBodyHeight = cardHeight - cardTitleHeight
            cardBackgroundShape.height = cardBodyHeight
        else:
            # Any card titles would be within the rest of the card
            cardBackgroundShape.top = card.top
            cardBodyHeight = cardHeight - cardTitleHeight
            cardBackgroundShape.height = cardHeight

        # Create any dividing line
        if (c > 0) & (cardShape == "line"):
            if cardLayout == "horizontal":
                # Dividing line is to the left of the card
                dividingLine = createLine(
                    cardBackgroundShape.left - int(Inches(horizontalCardGap / 2)),
                    cardBackgroundShape.top + Inches(0.75),
                    cardBackgroundShape.left - int(Inches(horizontalCardGap / 2)),
                    cardBackgroundShape.top + cardBackgroundShape.height - Inches(0.75),
                    slide.shapes,
                )
            else:
                # Dividing line is above the card
                dividingLine = createLine(
                    cardBackgroundShape.left + Inches(0.75),
                    cardBackgroundShape.top - int(Inches(verticalCardGap / 2)),
                    cardBackgroundShape.left
                    + cardBackgroundShape.width
                    - int(Inches(0.75)),
                    cardBackgroundShape.top - int(Inches(verticalCardGap / 2)),
                    slide.shapes,
                )

            # Perhaps set the line colour
            if cardDividerColour != ("None", ""):
                setColour(dividingLine.line.color, cardDividerColour)

            # Set the line width to a fixed 2pts
            dividingLine.line.width = Pt(2.0)

        # Position card body shape
        if (cardGraphicShape == None) & (cardMediaShape == None):
            # No graphic on this card
            cardBodyShape.top = card.bodyTop
            cardBodyShape.height = cardBodyHeight
            cardBodyShape.left = card.left
            cardBodyShape.width = cardWidth
        else:
            # Make room for graphic, audio, or video
            if cardGraphicPosition == "before":
                # Graphic before
                if cardLayout == "horizontal":
                    # Leave room above card body shape for graphic
                    cardBodyShape.top = card.bodyTop + cardGraphicSizeRaw + \
                        2 * cardGraphicPadding
                else:
                    # Don't leave room above card for graphic
                    cardBodyShape.top = card.bodyTop
                    cardBodyShape.width = cardWidth - 2 * cardGraphicPadding - \
                        cardGraphicSizeRaw

            else:
                # graphic after
                # Leave room below card body shape for graphic
                cardBodyShape.top = card.bodyTop

                if cardLayout == "vertical":
                    cardBodyShape.width = cardWidth - 2 * cardGraphicPadding - \
                        cardGraphicSizeRaw

            if cardLayout == "horizontal":
                # Calculate card body shape height, leaving room for any graphic
                cardBodyShape.height = cardBodyHeight - cardGraphicSizeRaw - \
                    2 * cardGraphicPadding
            else:
                cardBodyShape.height = cardBodyHeight

            # Scale graphic, audio, or video
            (cardMediaWidth, cardMediaHeight, scaledByHeight) = scalePicture(
                cardGraphicSizeRaw,
                cardGraphicSizeRaw,
                cardMediaNativeWidth,
                cardMediaNativeHeight,
            )

            if cardGraphicShape is not None:
                cardMediaShape = cardGraphicShape
            else:
                cardMediaShape = cardMediaShape

            # Vertically position graphic shape
            if cardGraphicPosition == "before":
                # Graphic before card text
                if cardLayout == "horizontal":
                    cardMediaShape.top = cardTitleShape.top + cardGraphicPadding + \
                        cardTitleShape.height + int((cardGraphicSizeRaw - cardMediaHeight) / 2)
                else:
                    cardMediaShape.top = cardBodyShape.top + \
                        int((cardBodyHeight - cardMediaHeight) / 2)
            else:
                # Graphic after card text
                if cardLayout == "horizontal":
                    cardMediaShape.top = cardBodyShape.height + cardBodyShape.top + \
                        cardGraphicPadding
                else:
                    cardMediaShape.top = cardBodyShape.top + \
                        int((cardBodyHeight - cardGraphicSizeRaw) / 2)

            # Horizontally position card graphic shape
            if cardLayout == "horizontal":
                cardMediaShape.left = cardTitleShape.left + \
                    int((cardTitleShape.width - cardMediaWidth) / 2)
            else:
                if cardGraphicPosition == "before":
                    cardMediaShape.left = card.left + cardGraphicPadding
                else:
                    cardMediaShape.left = card.left + cardBodyShape.width + \
                        cardGraphicPadding

            # Set dimensions of card graphic shape
            cardMediaShape.width = int(cardMediaWidth)
            cardMediaShape.height = int(cardMediaHeight)

        # Render any card body text
        if card.bullets != "":
            renderText(cardBodyShape, card.bullets)

        # Handle card background border line
        lf = cardBackgroundShape.line

        if (cardBorderColour != ("None", "")) & (cardShape != "line"):
            # Set border line colour
            setColour(lf.color, cardBorderColour)

        if cardShape == "line":
            # Lines between cards means cards have no border
            lf.fill.background()
        elif cardBorderWidth > 0:
            # Non-zero border line width
            lf.width = Pt(cardBorderWidth)
        elif cardBorderWidth == 0:
            # Zero border line width
            lf.fill.background()

        # Create any card shadow
        if cardShadow:
            createShadow(cardBackgroundShape)

    return slide


def createTableBlock(slideInfo, slide, renderingRectangle, tableBlockNumber):
    tableRows = slideInfo.tableRows[tableBlockNumber]
    tableMargin = md2pptx.globals.processingOptions.getCurrentOption("tableMargin")
    marginBase = md2pptx.globals.processingOptions.getCurrentOption("marginBase")
    baseTextSize = md2pptx.globals.processingOptions.getCurrentOption("baseTextSize")
    tableShadow = md2pptx.globals.processingOptions.getCurrentOption("tableShadow")

    printableTopLeftGraphicFilename = ""
    printableTopRightGraphicFilename = ""
    printablebottomLeftGraphicFilename = ""
    printableBottomRightGraphicFilename = ""

    # Handle table body
    if (len(tableRows) <= 2) & (len(tableRows[0]) <= 2):
        # This is a table with 1 or 2 rows and 1 or 2 columns
        isGraphicsGrid = True
        gridRows = len(tableRows)
        if gridRows == 1:
            gridColumns = len(tableRows[0])
        else:
            gridColumns = max(len(tableRows[0]), len(tableRows[1]))

        topGraphicCount = 0

        topLeftCellString = tableRows[0][0]
        # Attempt to retrieve media information for left side - top row
        (
            topLeftGraphicTitle,
            topLeftGraphicFilename,
            printableTopLeftGraphicFilename,
            topLeftGraphicHref,
            topLeftHTML,
            topLeftVideo,
            topGraphicCount,
        ) = parseMedia(topLeftCellString, topGraphicCount)

        # Attempt to retrieve filename for right side - top row
        if len(tableRows[0]) == 2:
            topRightCellString = tableRows[0][1]
        else:
            topRightCellString = ""

        (
            topRightGraphicTitle,
            topRightGraphicFilename,
            printableTopRightGraphicFilename,
            topRightGraphicHref,
            topRightHTML,
            topRightVideo,
            topGraphicCount,
        ) = parseMedia(topRightCellString, topGraphicCount)

        if topGraphicCount == 0:
            # Revert to normal table processing as no graphic spec in at least one cell
            isGraphicsGrid = False

        if gridRows == 2:
            # Attempt to retrieve filename for left side - bottom row
            bottomGraphicCount = 0

            bottomLeftCellString = tableRows[1][0]

            # Attempt to retrieve media information for left side - bottom row
            (
                bottomLeftGraphicTitle,
                bottomLeftGraphicFilename,
                printableBottomLeftGraphicFilename,
                bottomLeftGraphicHref,
                bottomLeftHTML,
                bottomLeftVideo,
                bottomGraphicCount,
            ) = parseMedia(bottomLeftCellString, bottomGraphicCount)

            # Attempt to retrieve filename for right side - bottom row
            if gridColumns == 2:
                if len(tableRows[1]) == 1:
                    # There is one cell in bottom row so this is centred "3-up"
                    bottomRightCellString = ""
                else:
                    bottomRightCellString = tableRows[1][1]
            else:
                bottomRightCellString = ""

            (
                bottomRightGraphicTitle,
                bottomRightGraphicFilename,
                printableBottomRightGraphicFilename,
                bottomRightGraphicHref,
                bottomRightHTML,
                bottomRightVideo,
                bottomGraphicCount,
            ) = parseMedia(bottomRightCellString, bottomGraphicCount)

            if bottomGraphicCount == 0:
                # Revert to normal table processing as no graphic spec in at least one cell
                isGraphicsGrid = False

    else:
        # This is a normal table because it has too many rows or columns to be a graphics grid
        isGraphicsGrid = False

    if isGraphicsGrid == True:

        ####################################################################
        # Print the media filenames                                        #
        ####################################################################
        if gridColumns == 2:
            # Doing 1- or 2-row side-by-side graphics slide
            reportGraphicFilenames(
                printableTopLeftGraphicFilename, printableTopRightGraphicFilename
            )
        else:
            # Doing 2 row, single column graphics slide
            reportGraphicFilenames(printableTopLeftGraphicFilename)

        if gridRows == 2:
            # Second row of filenames
            if gridColumns == 2:
                reportGraphicFilenames(
                    printableBottomLeftGraphicFilename,
                    printableBottomRightGraphicFilename,
                )
            else:
                reportGraphicFilenames(printableBottomLeftGraphicFilename)

        ####################################################################
        # Get the media dimensions                                         #
        ####################################################################
        if topLeftGraphicFilename != "":
            topLeftMediaWidth, topLeftMediaHeight = getGraphicDimensions(
                topLeftGraphicFilename
            )
            if topLeftMediaWidth == -1:
                if gridRows == 2:
                    print(
                        "Missing top left image file: "
                        + printableTopLeftGraphicFilename
                    )
                else:
                    print("Missing left image file: " + printableTopLeftGraphicFilename)

                return slide

        elif topLeftVideo is not None:
            (
                topLeftMediaWidth,
                topLeftMediaHeight,
                topLeftVideoType,
                topLeftVideoData,
            ) = getVideoInfo(topLeftVideo)

            if topLeftMediaWidth == -1:
                if gridRows == 2:
                    print(
                        "Missing top left video file: "
                        + printableTopLeftGraphicFilename
                    )
                else:
                    print("Missing left video file: " + printableTopLeftGraphicFilename)

                return slide

        if gridColumns == 2:
            # Get top right image dimensions
            if topRightGraphicFilename != "":
                topRightMediaWidth, topRightMediaHeight = getGraphicDimensions(
                    topRightGraphicFilename
                )
                if topRightMediaWidth == -1:
                    if gridRows == 2:
                        print(
                            "Missing top right image file: "
                            + printableTopRightGraphicFilename
                        )
                    else:
                        print(
                            "Missing right image file: "
                            + printableTopRightGraphicFilename
                        )

                    return slide

            elif topRightVideo is not None:
                (
                    topRightMediaWidth,
                    topRightMediaHeight,
                    topRightVideoType,
                    topRightVideoData,
                ) = getVideoInfo(topRightVideo)

                if topRightMediaWidth == -1:
                    if gridRows == 2:
                        print(
                            "Missing top right video file: "
                            + printableTopRightGraphicFilename
                        )
                    else:
                        print(
                            "Missing right video file: "
                            + printableTopRightGraphicFilename
                        )

                    return slide

        if gridRows == 2:
            # Get bottom left image dimensions
            if bottomLeftGraphicFilename != "":
                bottomLeftMediaWidth, bottomLeftMediaHeight = getGraphicDimensions(
                    bottomLeftGraphicFilename
                )
                if bottomLeftMediaWidth == -1:
                    print(
                        "Missing bottom left image file: "
                        + printableBottomLeftGraphicFilename
                    )
                    return slide

            elif bottomLeftVideo is not None:
                (
                    bottomLeftMediaWidth,
                    bottomLeftMediaHeight,
                    bottomLeftVideoType,
                    bottomLeftVideoData,
                ) = getVideoInfo(bottomLeftVideo)

                if bottomLeftMediaWidth == -1:
                    if gridRows == 2:
                        print(
                            "Missing bottom left video file: "
                            + printableBottomLeftGraphicFilename
                        )
                    else:
                        print(
                            "Missing left video file: "
                            + printableBottomLeftGraphicFilename
                        )

                    return slide

            if gridColumns == 2:
                # Get bottom right image dimensions
                if bottomRightGraphicFilename != "":
                    (
                        bottomRightMediaWidth,
                        bottomRightMediaHeight,
                    ) = getGraphicDimensions(bottomRightGraphicFilename)

                    if bottomRightMediaWidth == -1:
                        print(
                            "Missing bottom right image file: "
                            + printableBottomRightGraphicFilename
                        )

                        return slide

                elif bottomRightVideo is not None:
                    (
                        bottomRightMediaWidth,
                        bottomRightMediaHeight,
                        bottomRightVideoType,
                        bottomRightVideoData,
                    ) = getVideoInfo(bottomRightVideo)

                    if bottomRightMediaWidth == -1:
                        if gridRows == 2:
                            print(
                                "Missing bottom right video file: "
                                + printableBottomRightGraphicFilename
                            )
                        else:
                            print(
                                "Missing right video file: "
                                + printableBottomRightGraphicFilename
                            )

                        return slide

        # Calculate maximum picture height on slide
        maxPicHeight = renderingRectangle.height

        if gridRows == 2:
            # Adjusted if two rows
            maxPicHeight = maxPicHeight / 2 + Inches(0.2)

        # Calculate maximum picture width on slide
        maxPicWidth = renderingRectangle.width
        if gridColumns == 2:
            # Adjusted if two columns
            maxPicWidth = maxPicWidth / 2 - marginBase

        # Calculate horizontal middle of graphics space
        midGraphicsSpaceX = renderingRectangle.left + renderingRectangle.width / 2

        ####################################################################
        # Calculate the size of each graphic - scaled by the above rect    #
        ####################################################################

        if (topLeftGraphicFilename != "") | (topLeftVideo is not None):
            (
                topLeftPicWidth,
                topLeftPicHeight,
                usingHeightToScale,
            ) = scalePicture(
                maxPicWidth, maxPicHeight, topLeftMediaWidth, topLeftMediaHeight
            )

            if usingHeightToScale:
                # Calculate horizontal start
                if (gridColumns == 2) and (
                    (topRightGraphicFilename != "") | (topRightVideo is not None)
                ):
                    # Align top left media item to the left
                    topLeftPicX = (
                        renderingRectangle.left
                        + (midGraphicsSpaceX - marginBase - topLeftPicWidth) / 2
                    )
                else:
                    # Center sole top media item
                    topLeftPicX = midGraphicsSpaceX - topLeftPicWidth / 2
            else:
                # Calculate horizontal start
                if (gridColumns == 2) and (
                    (topRightGraphicFilename != "") | (topRightVideo is not None)
                ):
                    # Align top left media item to the left
                    topLeftPicX = renderingRectangle.left
                else:
                    # Center sole top media item
                    topLeftPicX = midGraphicsSpaceX - topLeftPicWidth / 2

            # Calculate vertical start
            topLeftPicY = renderingRectangle.top + (maxPicHeight - topLeftPicHeight) / 2

            if gridRows == 2:
                topLeftPicY -= Inches(0.2)

        if topLeftGraphicFilename != "":
            topLeftPicture = slide.shapes.add_picture(
                topLeftGraphicFilename,
                topLeftPicX,
                topLeftPicY,
                topLeftPicWidth,
                topLeftPicHeight,
            )

            if topLeftGraphicHref == "":
                topLeftGraphicHref = None

            pictureInfos.append(
                (topLeftPicture, topLeftGraphicHref, topLeftGraphicTitle)
            )
        elif topLeftVideo is not None:
            if topLeftVideoType == "Local":
                # Can use local file directly
                topLeftVideoShape = slide.shapes.add_movie(
                    topLeftVideo.source,
                    topLeftPicX,
                    topLeftPicY,
                    topLeftPicWidth,
                    topLeftPicHeight,
                    topLeftVideo.poster,
                )
            else:
                # First copy video data to temporary file
                tempVideoFile = tempfile.NamedTemporaryFile(
                    delete=False, suffix="mp4", dir=tempDir
                )
                tempVideoFile.write(topLeftVideoData)
                convertibleFilename = tempVideoFile.name
                tempVideoFile.close()

                # Use temporary file to make video
                topLeftVideo = slide.shapes.add_movie(
                    convertibleFilename,
                    topLeftPicX,
                    topLeftPicY,
                    topLeftPicWidth,
                    topLeftPicHeight,
                    topLeftVideo.poster,
                )

        if gridColumns == 2:
            # Top right media item
            if (topRightGraphicFilename != "") | (topRightVideo is not None):
                (
                    topRightPicWidth,
                    topRightPicHeight,
                    usingHeightToScale,
                ) = scalePicture(
                    maxPicWidth, maxPicHeight, topRightMediaWidth, topRightMediaHeight
                )

                if usingHeightToScale:
                    # Calculate horizontal start
                    topRightPicX = (
                        renderingRectangle.width + midGraphicsSpaceX - topRightPicWidth
                    ) / 2
                else:
                    # Calculate horizontal start
                    topRightPicX = (
                        renderingRectangle.width + midGraphicsSpaceX - topRightPicWidth
                    ) / 2

                # Calculate vertical start
                topRightPicY = (
                    renderingRectangle.top + (maxPicHeight - topRightPicHeight) / 2
                )

                if gridRows == 2:
                    topRightPicY -= Inches(0.2)

            if topRightGraphicFilename != "":
                topRightPicture = slide.shapes.add_picture(
                    topRightGraphicFilename,
                    topRightPicX,
                    topRightPicY,
                    topRightPicWidth,
                    topRightPicHeight,
                )

                if topRightGraphicHref == "":
                    topRightGraphicHref = None

                pictureInfos.append(
                    (topRightPicture, topRightGraphicHref, topRightGraphicTitle)
                )

            elif topRightVideo is not None:
                if topRightVideoType == "Local":
                    # Can use local file directly
                    topRightVideoShape = slide.shapes.add_movie(
                        topRightVideo.source,
                        topRightPicX,
                        topRightPicY,
                        topRightPicWidth,
                        topRightPicHeight,
                        topRightVideo.poster,
                    )
                else:
                    # First copy video data to temporary file
                    tempVideoFile = tempfile.NamedTemporaryFile(
                        delete=False, suffix="mp4", dir=tempDir
                    )
                    tempVideoFile.write(topRightVideoData)
                    convertibleFilename = tempVideoFile.name
                    tempVideoFile.close()

                    # Use temporary file to make video
                    topRightVideo = slide.shapes.add_movie(
                        convertibleFilename,
                        topRightPicX,
                        topRightPicY,
                        topRightPicWidth,
                        topRightPicHeight,
                        topRightVideo.poster,
                    )

        if gridRows == 2:
            # Need second row of media items
            # Bottom left media item
            if (bottomLeftGraphicFilename != "") | (bottomLeftVideo is not None):
                (
                    bottomLeftPicWidth,
                    bottomLeftPicHeight,
                    usingHeightToScale,
                ) = scalePicture(
                    maxPicWidth,
                    maxPicHeight,
                    bottomLeftMediaWidth,
                    bottomLeftMediaHeight,
                )

                if usingHeightToScale:
                    # Calculate horizontal start
                    if (gridColumns == 2) & (
                        (bottomRightGraphicFilename != "")
                        | (bottomRightVideo is not None)
                    ):
                        bottomLeftPicX = (
                            marginBase
                            + (midGraphicsSpaceX - marginBase - bottomLeftPicWidth) / 2
                        )
                    else:
                        bottomLeftPicX = midGraphicsSpaceX - bottomLeftPicWidth / 2
                else:
                    # Calculate horizontal start
                    if (gridColumns == 2) and (bottomRightGraphicFilename != ""):
                        # Align bottom left picture to the left
                        bottomLeftPicX = marginBase
                    else:
                        # Center sole bottom media item
                        bottomLeftPicX = midGraphicsSpaceX - bottomLeftPicWidth / 2

                # Calculate vertical start
                bottomLeftPicY = (
                    renderingRectangle.top + (maxPicHeight + bottomLeftPicHeight) / 2
                )

                if gridRows == 2:
                    bottomLeftPicY -= Inches(0.2)

            if bottomLeftGraphicFilename != "":
                bottomLeftPicture = slide.shapes.add_picture(
                    bottomLeftGraphicFilename,
                    bottomLeftPicX,
                    bottomLeftPicY,
                    bottomLeftPicWidth,
                    bottomLeftPicHeight,
                )

                if bottomLeftGraphicHref == "":
                    bottomLeftGraphicHref = None

                pictureInfos.append(
                    (bottomLeftPicture, bottomLeftGraphicHref, bottomLeftGraphicTitle)
                )

            elif bottomLeftVideo is not None:
                if bottomLeftVideoType == "Local":
                    # Can use local file directly
                    bottomLeftVideoShape = slide.shapes.add_movie(
                        bottomLeftVideo.source,
                        bottomLeftPicX,
                        bottomLeftPicY,
                        bottomLeftPicWidth,
                        bottomLeftPicHeight,
                        bottomLeftVideo.poster,
                    )
                else:
                    # First copy video data to temporary file
                    tempVideoFile = tempfile.NamedTemporaryFile(
                        delete=False, suffix="mp4", dir=tempDir
                    )
                    tempVideoFile.write(bottomLeftVideoData)
                    convertibleFilename = tempVideoFile.name
                    tempVideoFile.close()

                    # Use temporary file to make video
                    bottomLeftVideo = slide.shapes.add_movie(
                        convertibleFilename,
                        bottomLeftPicX,
                        bottomLeftPicY,
                        bottomLeftPicWidth,
                        bottomLeftPicHeight,
                        bottomLeftVideo.poster,
                    )

            if gridColumns == 2:
                # Bottom right media item
                if (bottomRightGraphicFilename != "") | (bottomRightVideo is not None):
                    (
                        bottomRightPicWidth,
                        bottomRightPicHeight,
                        usingHeightToScale,
                    ) = scalePicture(
                        maxPicWidth,
                        maxPicHeight,
                        bottomRightMediaWidth,
                        bottomRightMediaHeight,
                    )

                    if usingHeightToScale:
                        # Calculate horizontal start
                        bottomRightPicX = (
                            renderingRectangle.width
                            + midGraphicsSpaceX
                            - bottomRightPicWidth
                        ) / 2

                    else:
                        # Use the width to scale
                        # Calculate horizontal start
                        bottomRightPicX = (
                            renderingRectangle.width
                            + midGraphicsSpaceX
                            - bottomRightPicWidth
                        ) / 2

                    # Calculate vertical start
                    bottomRightPicY = (
                        renderingRectangle.top
                        + (maxPicHeight + bottomRightPicHeight) / 2
                    )

                    if gridRows == 2:
                        bottomRightPicY -= Inches(0.2)

                if bottomRightGraphicFilename != "":
                    if bottomRightGraphicFilename != "":
                        bottomRightPicture = slide.shapes.add_picture(
                            bottomRightGraphicFilename,
                            bottomRightPicX,
                            bottomRightPicY,
                            bottomRightPicWidth,
                            bottomRightPicHeight,
                        )

                        if bottomRightGraphicHref == "":
                            bottomRightGraphicHref = None

                        pictureInfos.append(
                            (
                                bottomRightPicture,
                                bottomRightGraphicHref,
                                bottomRightGraphicTitle,
                            )
                        )
                elif bottomRightVideo is not None:
                    if bottomRightVideoType == "Local":
                        # Can use local file directly
                        bottomRightVideoShape = slide.shapes.add_movie(
                            bottomRightVideo.source,
                            bottomRightPicX,
                            bottomRightPicY,
                            bottomRightPicWidth,
                            bottomRightPicHeight,
                            bottomRightVideo.poster,
                        )
                    else:
                        # First copy video data to temporary file
                        tempVideoFile = tempfile.NamedTemporaryFile(
                            delete=False, suffix="mp4", dir=tempDir
                        )
                        tempVideoFile.write(bottomRightVideoData)
                        convertibleFilename = tempVideoFile.name
                        tempVideoFile.close()

                        # Use temporary file to make video
                        bottomRightVideo = slide.shapes.add_movie(
                            convertibleFilename,
                            bottomRightPicX,
                            bottomRightPicY,
                            bottomRightPicWidth,
                            bottomRightPicHeight,
                            bottomRightVideo.poster,
                        )

    else:
        ################
        #              #
        # Normal table #
        #              #
        ################

        # Calculate maximum number of columns - as this is how wide we'll make the table
        columns = 0
        for row in tableRows:
            columns = max(columns, len(row))

        alignments = []
        widths = []

        # Adjust table if it contains a dash line as it's second line
        if len(tableRows) > 1:
            firstCellSecondRow = tableRows[1][0]
            if (firstCellSecondRow.startswith("-")) | (
                firstCellSecondRow.startswith(":-")
            ):
                haveTableHeading = True
            else:
                haveTableHeading = False
        else:
            haveTableHeading = False

        if haveTableHeading is True:
            # Has table heading
            tableHeadingBlurb = " with heading"

            # Figure out alignments of cells
            for cell in tableRows[1]:
                if cell.startswith(":-"):
                    if cell.endswith("-:"):
                        alignments.append("c")
                    else:
                        alignments.append("l")
                elif cell.endswith("-:"):
                    alignments.append("r")
                else:
                    alignments.append("l")

                widths.append(cell.count("-"))

            # Default any missing columns to left / single width
            if len(tableRows[1]) < columns:
                for _ in range(columns - len(tableRows[1])):
                    alignments.append("l")
                    widths.append(1)

            widths_total = sum(widths)

            # Remove this alignment / widths row from the table
            del tableRows[1]
        else:
            # No table heading
            tableHeadingBlurb = " without heading"

            # Use default width - 1 - and default alignment - l
            for c in range(columns):
                widths.append(1)
                alignments.append("l")

            # We don't know the widths so treat all equal
            widths_total = columns

        # Calculate number of rows
        rows = len(tableRows)
        alignments_count = len(alignments)

        # Create the table with the above number of rows and columns
        newTableShape = slide.shapes.add_table(rows, columns, 0, 0, 0, 0)

        newTable = newTableShape.table

        newTableShape.top = renderingRectangle.top
        newTableShape.left = renderingRectangle.left + tableMargin - marginBase
        newTableShape.height = min(renderingRectangle.height, Inches(0.25) * rows)
        newTableShape.width = renderingRectangle.width - 2 * (tableMargin - marginBase)
        shapeWidth = newTableShape.width

        # Perhaps create a drop shadow for a table
        if tableShadow:
            createShadow(newTable)

        # Set whether first row is not special
        newTable.first_row = haveTableHeading

        print(
            "           --> "
            + str(rows)
            + " x "
            + str(columns)
            + " table"
            + tableHeadingBlurb
        )

        # Set column widths
        cols = newTable.columns
        for colno in range(columns):
            cols[colno].width = int(shapeWidth * widths[colno] / widths_total)

        # Get options for filling in the cells
        compactTables = md2pptx.globals.processingOptions.getCurrentOption("compactTables")
        spanCells = md2pptx.globals.processingOptions.getCurrentOption("spanCells")
        tableHeadingSize = md2pptx.globals.processingOptions.getCurrentOption("tableHeadingSize")

        # Fill in the cells
        for rowNumber, row in enumerate(tableRows):
            # Add dummy cells to the end of the row so that there are as many
            # cells in the row as there are columns in the table
            cellCount = len(row)

            # Unless there is a non-empty cell there is no anchor cell for this row
            if spanCells == "yes":
                potentialAnchorCell = None

            for c in range(cellCount, columns):
                row.append("")

            for columnNumber, cell in enumerate(row):
                newCell = newTable.cell(rowNumber, columnNumber)

                if spanCells == "yes":
                    if cell != "":
                        potentialAnchorCell = newCell
                    else:
                        if potentialAnchorCell is not None:
                            # Might need to remove previous cell merge
                            if potentialAnchorCell.span_width > 1:
                                potentialAnchorCell.split()

                            # Merge the cells from the anchor up to this one
                            potentialAnchorCell.merge(newCell)

                # For compact table remove the margins around the text
                if compactTables > 0:
                    newCell.margin_top = Pt(0)
                    newCell.margin_bottom = Pt(0)

                newCell.text = ""
                text_frame = newCell.text_frame

                # Set cell's text alignment
                p = text_frame.paragraphs[0]

                # Set cell's text size - if necessary
                if baseTextSize > 0:
                    p.font.size = Pt(baseTextSize)

                # For compact table use specified point size for text
                if compactTables > 0:
                    p.font.size = Pt(compactTables)

                if (rowNumber == 0) & (tableHeadingSize > 0):
                    p.font.size = Pt(tableHeadingSize)

                if columnNumber >= alignments_count:
                    p.alignment = PP_ALIGN.LEFT
                elif alignments[columnNumber] == "r":
                    p.alignment = PP_ALIGN.RIGHT
                elif alignments[columnNumber] == "c":
                    p.alignment = PP_ALIGN.CENTER
                else:
                    p.alignment = PP_ALIGN.LEFT

                addFormattedText(p, cell)

        # Apply table border styling - whether there is any or not
        applyTableLineStyling(
            newTable,
            md2pptx.globals.processingOptions,
        )

    return slide


def createChevron(
    text,
    x,
    y,
    width,
    height,
    filled,
    shapes,
    fontSize,
    wantLink,
    unhighlightedBackground,
):
    global TOCruns

    # Create shape
    shape = shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, width, height)

    # Set shape's text
    shape.text = text

    # Set shape's text attributes
    tf = shape.text_frame
    p = tf.paragraphs[0]
    f = p.font
    f.size = Pt(fontSize)
    f.color.rgb = RGBColor(0, 0, 0)

    # If want link create it from the first run
    if wantLink:
        TOCruns.append(p.runs[0])

    # Set shape's outline attributes
    shape.line.color.rgb = RGBColor(0, 0, 0)
    shape.line.width = Pt(1.0)

    # Potentially fill background
    if filled is False:
        shape.fill.background()
    else:
        if wantLink & (unhighlightedBackground != ""):
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(unhighlightedBackground)


def createOval(
    text,
    x,
    y,
    width,
    height,
    filled,
    shapes,
    fontSize,
    wantLink,
    unhighlightedBackground,
):
    global TOCruns

    # Create shape
    shape = shapes.add_shape(MSO_SHAPE.OVAL, x, y, width, height)

    # Set shape's text
    shape.text = text

    # Set shape's text attributes
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    f = p.font
    f.size = Pt(fontSize)
    f.color.rgb = RGBColor(0, 0, 0)

    # If want link create it from the first run
    if wantLink:
        TOCruns.append(p.runs[0])

    # Set shape's outline attributes
    shape.line.color.rgb = RGBColor(191, 191, 191)
    shape.line.width = Pt(1.0)

    # Potentially fill background
    if filled is False:
        shape.fill.background()
        shape.line.width = Pt(3.0)
    else:
        if wantLink & (unhighlightedBackground != ""):
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(unhighlightedBackground)


def createLine(x0, y0, x1, y1, shapes, colour=("RGB", "#BFBFBF"), width=4.0):
    # Create line
    line = shapes.add_shape(MSO_SHAPE.LINE_INVERSE, x0, y0, x1 - x0, y1 - y0)

    # Set shape's outline attributes
    setColour(line.line.color, colour)

    line.line.width = Pt(width)

    return line


def delinkify(text):
    if linkMatch := linkRegex.match(text):
        linkText = linkMatch.group(1)
        linkURL = linkMatch.group(2)
        return (linkText, linkURL)

    elif linkMatch := indirectReferenceUsageRegex(text):
        print(linkMatch.group(1))
        print(linkMatch.group(2))
        return (text, "")

    else:
        return (text, "")


def createTOCSlide(presentation, slideNumber, titleText, bullets, tocStyle):
    global SectionSlides
    titleOnlyLayout = md2pptx.globals.processingOptions.getCurrentOption("titleOnlyLayout")
    blankLayout = md2pptx.globals.processingOptions.getCurrentOption("blankLayout")
    tocTitle = md2pptx.globals.processingOptions.getCurrentOption("tocTitle")
    marginBase = md2pptx.globals.processingOptions.getCurrentOption("marginBase")
    pageTitleSize = md2pptx.globals.processingOptions.getCurrentOption("pageTitleSize")
    pageSubtitleSize = md2pptx.globals.processingOptions.getCurrentOption("pageSubtitleSize")

    if tocStyle != "plain":
        if titleText == tocTitle:
            reportSlideTitle(
                slideNumber, 3, f'Table Of Contents (Style: "{tocStyle}") {titleText}'
            )

        else:
            reportSlideTitle(slideNumber, 2, titleText)

    if tocStyle == "plain":
        if titleText != tocTitle:
            slide = createTitleOrSectionSlide(
                presentation,
                slideNumber,
                titleText,
                md2pptx.globals.processingOptions.getCurrentOption("sectionSlideLayout"),
                md2pptx.globals.processingOptions.getCurrentOption("sectionTitleSize"),
                slideInfo.subtitleText,
                md2pptx.globals.processingOptions.getCurrentOption("sectionSubtitleSize"),
                notes_text,
            )
        else:
            # Remove the links from the bullets and replace with target slide title
            for bullet in bullets:
                linkMatch = linkRegex.match(bullet[1])
                bullet[1] = linkMatch.group(1)

            # Create the TOC slide - with these neutralised titles
            slide = createContentSlide(
                presentation,
                slideNumber,
                slideInfo,
            )

            # Postprocess slide to pick up runs - for TOC creation
            body = findBodyShape(slide)
            text_frame = body.text_frame
            for p in text_frame.paragraphs:
                TOCruns.append(p.runs[0])

        # Store the new slide in the list of section slides - for fixing up links
        SectionSlides[titleText] = slide

        return slide

    else:
        slide = addSlide(
            presentation, presentation.slide_layouts[titleOnlyLayout], None
        )
        title = findTitleShape(slide)

    SectionSlides[titleText] = slide

    shapes = slide.shapes

    # Add title if TOC slide. Or delete shape if not
    if titleText == tocTitle:
        # Is TOC slide so add title
        slideTitleBottom, title, flattenedTitle = formatTitle(
            presentation, slide, tocTitle, pageTitleSize, pageSubtitleSize
        )
    else:
        # Is not TOC slide so delete title shape and adjust where title bottom
        # would be
        deleteSimpleShape(title)
        slideTitleBottom = marginBase

    # Get the rectangle the content will draw in
    contentLeft, contentWidth, contentTop, contentHeight = getContentRect(
        presentation, slide, slideTitleBottom, marginBase
    )

    # Create global list of TOC entries
    for bullet in bullets:
        bulletLevel, bulletText, bulletType = bullet
        if bulletLevel == 0:
            # Level 0 is top level so create a TOC entry
            linkText, linkHref = delinkify(bulletText)
            TOCEntries.append([linkText, linkHref])

    TOCEntryCount = len(TOCEntries)

    TOCFontSize = md2pptx.globals.processingOptions.getCurrentOption("TOCFontSize")

    TOCItemHeight = md2pptx.globals.processingOptions.getCurrentOption("TOCItemHeight")

    TOCItemColour = md2pptx.globals.processingOptions.getCurrentOption("TOCItemColour")

    height = Inches(TOCItemHeight)

    if tocStyle == "chevron":
        if height == 0:
            height = Inches(1)

        width = height * 2.5

        entryGap = Inches(-0.5 * height / Inches(1))

        if TOCFontSize == 0:
            TOCFontSize = 14

    elif tocStyle == "circle":
        if height == 0:
            height = Inches(1.25)

        width = height

        entryGap = Inches(0.5)

        if TOCFontSize == 0:
            TOCFontSize = 12

    rowGap = Inches(md2pptx.globals.processingOptions.getCurrentOption("TOCRowGap"))

    TOCEntriesPerRow = int(
        (presentation.slide_width - 2 * marginBase) / (width + entryGap)
    )

    rowCount = 1 + TOCEntryCount / TOCEntriesPerRow

    # Calculate actual TOC height so it can be vertically centred
    TOCHeight = (rowCount * height) + ((rowCount - 1) * rowGap)

    # Calculate where top of TOC should be
    TOCtop = slideTitleBottom + (contentHeight - TOCHeight + height) / 2

    # Calculate actual TOC width
    TOCWidth = TOCEntriesPerRow * (width + entryGap)

    # Calculate where the TOC will start
    TOCleft = (presentation.slide_width - TOCWidth + entryGap) / 2

    x = TOCleft
    y = TOCtop

    AbsoluteTOCEntryNumber = 1

    TOCEntryNumber = 1

    for entry in TOCEntries:
        entryText = entry[0]
        entryHref = entry[1]

        if entryText == titleText:
            wantFilled = False
            wantLink = False
        else:
            wantFilled = True
            wantLink = True

        if tocStyle == "chevron":
            createChevron(
                entryText,
                x,
                y,
                width,
                height,
                wantFilled,
                shapes,
                TOCFontSize,
                wantLink,
                TOCItemColour,
            )

        elif tocStyle == "circle":
            # Create the circle
            createOval(
                entryText,
                x,
                y,
                width,
                height,
                wantFilled,
                shapes,
                TOCFontSize,
                wantLink,
                TOCItemColour,
            )

            # Create half connector to next one - if not last
            if AbsoluteTOCEntryNumber < TOCEntryCount:
                connector = createLine(
                    x + width,
                    y + height / 2,
                    x + width + entryGap / 2,
                    y + height / 2,
                    shapes,
                )

            # Create half connector to previous one - if not first
            if AbsoluteTOCEntryNumber > 1:
                # z =1
                connector = createLine(
                    x - entryGap / 2, y + height / 2, x, y + height / 2, shapes
                )

        # Prepare for the next TOC entry - even if there isn't one
        x = x + width + entryGap

        # If beyond end of line the next TOC entry would be at the start of the next line
        AbsoluteTOCEntryNumber = AbsoluteTOCEntryNumber + 1
        TOCEntryNumber = TOCEntryNumber + 1
        if TOCEntryNumber == TOCEntriesPerRow + 1:
            x = TOCleft
            y = y + rowGap + height
            TOCEntryNumber = 1

    if want_numbers_content is True:
        addFooter(presentation, slideNumber, slide)

    return slide


def createSlide(presentation, slideNumber, slideInfo):
    abstractTitle = md2pptx.globals.processingOptions.getCurrentOption("abstractTitle")
    tocTitle = md2pptx.globals.processingOptions.getCurrentOption("tocTitle")
    tocStyle = md2pptx.globals.processingOptions.getCurrentOption("tocStyle")
    sectionTitleSize = md2pptx.globals.processingOptions.getCurrentOption("sectionTitleSize")
    presTitleSize = md2pptx.globals.processingOptions.getCurrentOption("presTitleSize")
    sectionSubtitleSize = md2pptx.globals.processingOptions.getCurrentOption("sectionSubtitleSize")
    presSubtitleSize = md2pptx.globals.processingOptions.getCurrentOption("presSubtitleSize")
    leftFooterText = md2pptx.globals.processingOptions.getCurrentOption("leftFooterText")
    footerfontsizespec = md2pptx.globals.processingOptions.getCurrentOption("footerFontSize")
    middleFooterText = md2pptx.globals.processingOptions.getCurrentOption("middleFooterText")
    rightFooterText = md2pptx.globals.processingOptions.getCurrentOption("rightFooterText")
    sectionFooters = md2pptx.globals.processingOptions.getCurrentOption("sectionFooters")
    liveFooters = md2pptx.globals.processingOptions.getCurrentOption("liveFooters")
    transition = md2pptx.globals.processingOptions.getCurrentOption("transition")
    hidden = md2pptx.globals.processingOptions.getCurrentOption("hidden")

    if slideInfo.blockType in ["content", "code", "table"]:
        if (tocStyle != "") & (tocTitle == slideInfo.titleText):
            # This is a Table Of Contents slide
            slide = createTOCSlide(
                presentation,
                slideNumber,
                slideInfo.titleText,
                slideInfo.bullets,
                tocStyle,
            )
        elif (abstractTitle != "") & (abstractTitle == slideInfo.titleText):
            # This is an abstract slide
            slide = createAbstractSlide(
                presentation,
                slideNumber,
                slideInfo.titleText,
                slideInfo.bullets,
            )
        else:
            # This is an ordinary contents slide
            slide = createContentSlide(
                presentation,
                slideNumber,
                slideInfo,
            )

    elif slideInfo.blockType == "section":
        if tocStyle != "":
            # This is a section slide in TOC style
            slide = createTOCSlide(
                presentation,
                slideNumber,
                slideInfo.titleText,
                slideInfo.bullets,
                tocStyle,
            )
        else:
            slide = createTitleOrSectionSlide(
                presentation,
                slideNumber,
                slideInfo.titleText,
                md2pptx.globals.processingOptions.getCurrentOption("sectionSlideLayout"),
                sectionTitleSize,
                slideInfo.subtitleText,
                sectionSubtitleSize,
                notes_text,
            )

    elif slideInfo.blockType == "title":
        slide = createTitleOrSectionSlide(
            presentation,
            slideNumber,
            slideInfo.titleText,
            md2pptx.globals.processingOptions.getCurrentOption("titleSlideLayout"),
            presTitleSize,
            slideInfo.subtitleText,
            presSubtitleSize,
            notes_text,
        )

    if footerfontsizespec == "":
        footerFontSize = Pt(8.0)
    else:
        footerFontSize = Pt(footerfontsizespec)

    footerBoxTop = prs.slide_height - numbersHeight / 2 - footerFontSize
    footerBoxHeight = footerFontSize * 2

    if slideInfo.blockType in ["title", "section"]:
        if sectionFooters == "yes":
            wantFooters = True
        else:
            wantFooters = False

        if slideInfo.blockType == "section":
            prs.lastSectionTitle = slideInfo.titleText.strip()
            prs.lastSectionSlide = slide
        elif slideInfo.blockType == "title":
            prs.lastPresTitle = slideInfo.titleText.strip()
            prs.lastPresSubtitle = slideInfo.subtitleText.strip()

    else:
        wantFooters = True

    if wantFooters:
        # Left pseudo-footer
        if leftFooterText != "":
            leftFooterMargin = Inches(0.5)
            leftFooterBoxLeft = leftFooterMargin
            leftFooterBoxWidth = prs.slide_width / 3 - leftFooterMargin
            leftFooter = slide.shapes.add_textbox(
                leftFooterBoxLeft, footerBoxTop, leftFooterBoxWidth, footerBoxHeight
            )

            leftFooter.text, wantHyperLink = substituteFooterVariables(
                leftFooterText, liveFooters
            )

            if wantHyperLink:
                createShapeHyperlinkAndTooltip(leftFooter, prs.lastSectionSlide, "")

            for fp in leftFooter.text_frame.paragraphs:
                fp.alignment = PP_ALIGN.LEFT
                fp.font.size = footerFontSize

        # Middle pseudo-footer
        if middleFooterText != "":
            middleFooterBoxLeft = prs.slide_width / 3
            middleFooterBoxWidth = prs.slide_width / 3
            middleFooter = slide.shapes.add_textbox(
                middleFooterBoxLeft, footerBoxTop, middleFooterBoxWidth, footerBoxHeight
            )

            middleFooter.text, wantHyperLink = substituteFooterVariables(
                middleFooterText, liveFooters
            )

            if wantHyperLink:
                createShapeHyperlinkAndTooltip(middleFooter, prs.lastSectionSlide, "")

            for fp in middleFooter.text_frame.paragraphs:
                fp.alignment = PP_ALIGN.CENTER
                fp.font.size = footerFontSize

        # Right pseudo-footer
        if rightFooterText != "":
            rightFooterMargin = Inches(0.25)
            rightFooterBoxLeft = prs.slide_width * 2 / 3
            rightFooterBoxWidth = prs.slide_width / 3 - rightFooterMargin
            rightFooter = slide.shapes.add_textbox(
                rightFooterBoxLeft, footerBoxTop, rightFooterBoxWidth, footerBoxHeight
            )

            rightFooter.text, wantHyperLink = substituteFooterVariables(
                rightFooterText, liveFooters
            )

            if wantHyperLink:
                createShapeHyperlinkAndTooltip(rightFooter, prs.lastSectionSlide, "")

            for fp in rightFooter.text_frame.paragraphs:
                fp.alignment = PP_ALIGN.RIGHT
                fp.font.size = footerFontSize

    slideNumber = slideNumber + 1

    sequence = []

    addSlideTransition(slide, transition)

    if hidden:
        slide._element.set('show', '0')
    else:
        slide._element.set('show', '1')

    return [slideNumber, slide, sequence]


# Add a transition effect - for transitioning INTO the slide.
def addSlideTransition(slide, transitionType):
    # Handle "no transition" case
    if (transitionType == "none") | (transitionType == ""):
        return

    if transitionType in [
        "fracture",
    ]:
        choiceNS = "p15"
    else:
        choiceNS = "p14"

    # Construct first boilerplate XML fragment
    xml = '      <mc:AlternateContent ' + namespacesFragment(["mc"]) + '>\n'
    xml += (
        "    <mc:Choice xmlns:"
        + choiceNS
        + '="'
        + namespaceURL[choiceNS]
        + '" Requires="'
        + choiceNS
        + '">\n'
    )
    xml += (
        '<p:transition ' + namespacesFragment(["p", "p14"]) + 'spd="slow" p14:dur="3400">\n'
    )

    # Add in transition element
    if transitionType in [
        "wipe",
    ]:
        xml += "         <p:" + transitionType + " />\n"

    elif transitionType in [
        "push",
    ]:
        xml += "         <p:" + transitionType + ' dir="u"/>\n'

    elif transitionType in [
        "vortex",
    ]:
        xml += "         <p14:" + transitionType + ' dir="r"/>\n'

    elif transitionType in [
        "split",
    ]:
        xml += "         <p:" + transitionType + ' orient="vert"/>\n'

    elif transitionType in [
        "fracture",
    ]:
        xml += '                 <p15:prstTrans prst="fracture" />\n'

    else:
        xml += "         <p14:" + transitionType + " />\n"

    # Construct last boilerplate XML fragment

    xml += """
      </p:transition>
    </mc:Choice>
    <mc:Fallback>
    """

    xml += '      <p:transition ' + namespacesFragment(["p"]) + ' spd="slow">\n'

    if transitionType in [
        "split",
    ]:
        xml += (
            "        <p:"
            + transitionType
            + ' orient="vert" ' + namespacesFragment(["p"]) + '/>\n'
        )

    else:
        xml += "        <p:fade />\n"

    xml += """
      </p:transition>
    </mc:Fallback>
  </mc:AlternateContent>
   """

    # Turn this into an XML fragment
    xmlFragment = parse_xml(xml)

    # Add to slide's XML
    slide.element.insert(-1, xmlFragment)


def createTaskSlides(prs, slideNumber, tasks, titleStem):
    tasksPerPage = md2pptx.globals.processingOptions.getCurrentOption("tasksPerPage")

    taskSlideNumber = 0

    taskCount = len(tasks)
    for taskNumber, task in enumerate(tasks):
        if taskNumber % tasksPerPage == 0:
            # Is first task in a page
            if taskNumber > 0:
                # Print a "tasks" slide - as we have one to print out
                taskSlideNumber += 1
                if taskCount > tasksPerPage:
                    # More than one task page
                    title = titleStem + " - " + str(taskSlideNumber)
                else:
                    # Only one task page
                    title = titleStem

                taskBlock = [taskRows]

                slideInfo = SlideInfo(
                    title, "", "table", [], taskBlock, [], [], ["table"]
                )
                slide = createContentSlide(prs, slideNumber, slideInfo)

                # Fix up references to be active links to the slide where the task
                # was declared
                table = findBodyShape(slide).table
                for row in table.rows:
                    cell0Text = row.cells[0].text
                    if cell0Text not in ["Slide", ""]:
                        # First cell refers to a specific slide number - so link to it
                        run = row.cells[0].text_frame.paragraphs[0].runs[0]
                        createRunHyperlinkOrTooltip(
                            run, prs.slides[int(cell0Text) - 2 + templateSlideCount], ""
                        )

                slideNumber += 1

            taskRows = [["Slide", "Due", "Task", "Tags", "Done"]]
            taskRows.append(["-:", ":--:", ":----", ":----", ":--:"])
            old_sNum = 0

        sNum, taskText, dueDate, tags, done = task

        if tags != "":
            # Sort tags - if there are any
            tagList = re.split("[, ]", tags)
            sortedTagList = sorted(tagList)
            tags = str.join(",", sortedTagList)

        if sNum != old_sNum:
            taskRows.append([str(sNum), dueDate, taskText, tags, done])
        else:
            taskRows.append(["", dueDate, taskText, tags, done])
        old_sNum = sNum

    # Print a final "tasks" slide
    taskSlideNumber += 1
    if taskCount > tasksPerPage:
        title = titleStem + " - " + str(taskSlideNumber)
    else:
        title = titleStem

    taskBlock = [taskRows]
    slideInfo = SlideInfo(title, "", "table", [], taskBlock, [], [], ["table"])
    slide = createContentSlide(prs, slideNumber, slideInfo)

    # Fix up references to be active links to the slide where the task
    # was declared
    table = findBodyShape(slide).table
    for row in table.rows:
        cell0Text = row.cells[0].text
        if cell0Text not in ["Slide", ""]:
            # First cell refers to a specific slide number - so link to it
            run = row.cells[0].text_frame.paragraphs[0].runs[0]
            createRunHyperlinkOrTooltip(
                run, prs.slides[int(cell0Text) - 2 + templateSlideCount], ""
            )

    slideNumber += 1


def createGlossarySlides(prs, slideNumber, abbrevDictionary):
    termSlideNumber = 0
    glossarySlides = []

    glossaryTitle = md2pptx.globals.processingOptions.getCurrentOption("glossaryTitle")
    glossaryTerm = md2pptx.globals.processingOptions.getCurrentOption("glossaryTerm")
    glossaryTermsPerPage = md2pptx.globals.processingOptions.getCurrentOption("glossaryTermsPerPage")
    glossaryMeaningWidth = md2pptx.globals.processingOptions.getCurrentOption("glossaryMeaningWidth")
    glossaryMeaning = md2pptx.globals.processingOptions.getCurrentOption("glossaryMeaning")

    termCount = len(abbrevDictionary)

    for termNumber, term in enumerate(sorted(abbrevDictionary.keys())):
        if termNumber % glossaryTermsPerPage == 0:
            # Is first glossary term in a page
            if termNumber > 0:
                # Print a "glossary" slide - as we have one to print out
                termSlideNumber += 1
                if termCount > glossaryTermsPerPage:
                    # More than one glossary page
                    title = glossaryTitle + " - " + str(termSlideNumber)
                else:
                    # Only one glossary page
                    title = glossaryTerm

                glossaryBlock = [glossaryRows]
                slideInfo = SlideInfo(
                    title, "", "table", [], glossaryBlock, [], [], ["table"]
                )
                slide = createContentSlide(prs, slideNumber, slideInfo)

                glossarySlides.append(slide)
                slideNumber += 1

            glossaryRows = [[glossaryTerm, glossaryMeaning]]
            glossaryRows.append([":-", ":" + ("-" * glossaryMeaningWidth)])
            old_sNum = 0

        meaning = abbrevDictionary.get(term)

        glossaryRows.append([term, meaning])

    # Print a final "glossary" slide
    termSlideNumber += 1
    if termCount > glossaryTermsPerPage:
        # More than one glossary page
        title = glossaryTitle + " - " + str(termSlideNumber)
    else:
        # Only one glossary page
        title = glossaryTitle

    glossaryBlock = [glossaryRows]
    slideInfo = SlideInfo(title, "", "table", [], glossaryBlock, [], [], ["table"])
    slide = createContentSlide(prs, slideNumber, slideInfo)
    glossarySlides.append(slide)
    slideNumber += 1

    return slideNumber, glossarySlides


def createSlideNotes(slide, notes_text):
    # Remove surrounding white space
    notes_text = notes_text.strip().lstrip("\n")

    if slide.notes_slide.notes_text_frame.text != "":
        # Notes already filled in
        return

    if notes_text != "":
        # There is substantive slide note text so create the note
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame

        # addFormattedText handles eg hyperlinks and entity references
        addFormattedText(text_frame.paragraphs[0], notes_text)


def createFootnoteSlides(prs, slideNumber, footnoteDefinitions):
    footnotesSlideNumber = 0
    footnoteSlides = []

    footnotesTitle = md2pptx.globals.processingOptions.getCurrentOption("footnotesTitle")
    footnotesPerPage = md2pptx.globals.processingOptions.getCurrentOption("footnotesPerPage")

    footnoteCount = len(footnoteDefinitions)

    for footnoteNumber, footnote in enumerate(footnoteDefinitions):
        if footnoteNumber % footnotesPerPage == 0:
            # Is first footnote in a page
            if footnoteNumber > 0:
                # Print a "footnotes" slide - as we have one to print out
                footnotesSlideNumber += 1
                if footnoteCount > footnotesPerPage:
                    # More than one footnotes page
                    title = footnotesTitle + " - " + str(footnotesSlideNumber)
                else:
                    # Only one footnotes page
                    title = footnotesTitle

                slideInfo = SlideInfo(
                    title, "", "content", bullets, [], cards, [], ["list"]
                )
                slideNumber, slide, sequence = createSlide(prs, slideNumber, slideInfo)

                footnoteSlides.append(slide)

                # Turn off bulleting
                removeBullets(findBodyShape(slide).text_frame)

                slideNumber += 1
            bullets = []
            old_sNum = 0

        bullets.append(
            [
                1,
                str(footnoteNumber + 1) + ". " + footnoteDefinitions[footnoteNumber][1],
                "bullet",
            ]
        )

    # Print a final "footnote" slide
    footnotesSlideNumber += 1
    if footnoteCount > footnotesPerPage:
        # More than one footnotes page
        title = footnotesTitle + " - " + str(footnotesSlideNumber)
    else:
        # Only one footnotes page
        title = footnotesTitle

    slideInfo = SlideInfo(title, "", "content", bullets, [], cards, [], ["list"])
    slideNumber, slide, sequence = createSlide(prs, slideNumber, slideInfo)

    footnoteSlides.append(slide)

    # Turn off bulleting
    removeBullets(findBodyShape(slide).text_frame)

    slideNumber += 1

    return slideNumber, footnoteSlides


def cli():
    start_time = time.time()

    banner = (
        "md2pptx Markdown To Powerpoint Converter " + md2pptx_level + " " + md2pptx_date
    )

    bannerUnderline = ""
    for i in range(len(banner)):
        bannerUnderline = bannerUnderline + "="

    print("\n" + banner + "\n" + bannerUnderline)
    print("\nOpen source project: https://github.com/MartinPacker/md2pptx")

    print("\nExternal Dependencies:")
    print("\n  Python: " + platform.python_version())

    print("  python-pptx: " + pptx_version)

    if have_pillow:
        print("  Pillow: " + PIL.__version__)
    else:
        print("  Pillow: Not Installed")

    if have_cairosvg:
        print("  CairoSVG: " + cairosvg.__version__)
    else:
        print("  CairoSVG: Not Installed")

    if have_graphviz:
        print("  graphviz: " + graphviz.__version__)
    else:
        print("  graphviz: Not Installed")

    print("\nInternal Dependencies:")
    print(f"\n  funnel: {md2pptx.funnel.version}")
    print(f"  runPython: {md2pptx.runPython.version}")

    input_file = []

    if len(sys.argv) > 2:
        # Have input file as well as output file
        input_filename = sys.argv[1]
        output_filename = sys.argv[2]

        if Path(input_filename).exists():
            input_path = Path(input_filename)

            with input_path.open(mode='r', encoding='utf-8') as file:
                input_file = file.readlines()
        else:
            print("Input file specified but does not exist. Terminating.")
    elif len(sys.argv) == 1:
        print("No parameters. Terminating")
        sys.exit()
    else:
        output_filename = sys.argv[1]

        input_file = sys.stdin.readlines()

    if len(input_file) == 0:
        print("Empty input file. Terminating")
        sys.exit()

    slideNumber = 1

    bulletRegex = re.compile(r"^(\s)*(\*)(.*)")
    numberRegex = re.compile(r"^(\s)*(\d+)\.(.*)")
    metadataRegex = re.compile("^(.+):(.+)")

    graphicRE = r"!\[(.*?)\]\((.+?)\)"
    graphicRegex = re.compile(graphicRE)

    clickableGraphicRE = r"\[" + graphicRE + r"\]\((.+?)\)"
    clickableGraphicRegex = re.compile(clickableGraphicRE)

    videoRE = "<video (.*?)></video>"
    videoRegex = re.compile(videoRE)

    audioRE = "<audio (.*?)></audio>"
    audioRegex = re.compile(audioRE)

    linkRegex = re.compile(r"^\[(.+)\]\((.+)\)")
    footnoteDefinitionRegex = re.compile(r"^\[\^(.+?)\]: (.+)")
    slideHrefRegex = re.compile(r"(.+)\[(.+)\]$")
    anchorRegex = re.compile("^<a id=[\"'](.+)[\"']></a>")
    dynamicMetadataRegex = re.compile("^<!-- md2pptx: (.+): (.+) -->")
    indirectReferenceAnchorRegex = re.compile(r"^\[(.+?)\]: (.+)")
    indirectReferenceUsageRegex = re.compile(r"[(.+?)]\[(.+?)]")

    # Default slide layout enumeration
    md2pptx.globals.processingOptions.setOptionValuesArray(
        [
            ["titleSlideLayout", 0],
            ["sectionSlideLayout", 1],
            ["contentSlideLayout", 2],
            ["titleOnlyLayout", 5],
            ["blanklayout", 6],
        ]
    )

    # Abbreviation Dictionary
    abbrevDictionary = {}

    # Abbreviation Runs Dictionary
    abbrevRunsDictionary = {}

    # Footnote runs Dictionary
    footnoteRunsDictionary = {}

    # Extract metadata
    metadata_lines = []
    afterMetadataAndHTML = []


    TOCruns = []
    SectionSlides = {}

    inMetadata = True
    in_comment = False
    inHTML = False
    inCode = False

    # Pass 1: Strip out comments and metadata, storing the latter
    for line in input_file:
        if line.lstrip().startswith("<!-- md2pptx: "):
            # md2pptx dynamic metadata so keep it
            afterMetadataAndHTML.append(line)

        if line.lstrip().startswith("<!--"):
            if line.rstrip().endswith("-->"):
                # Note: Not taking text after end of comment
                continue
            else:
                in_comment = True
                continue

        elif line.rstrip().endswith("-->"):
            # Note: Not taking text after end of comment
            in_comment = False
            continue

        elif in_comment is True:
            continue

        elif (line.lstrip()[:1] == "<") & (inCode is False):
            lineLstrip = line.lstrip()
            if startswithOneOf(lineLstrip, ["<a id=", "<span "]):
                # Line goes to post-metadata array
                afterMetadataAndHTML.append(line)

            elif startswithOneOf(lineLstrip, ["<code>", "<pre>"]):
                inCode = True
                afterMetadataAndHTML.append(line)

            elif startswithOneOf(lineLstrip, ["</code>", "</pre>"]):
                inCode = False
                afterMetadataAndHTML.append(line)

            elif startswithOneOf(lineLstrip, ["<video ", "<audio "]):
                # Line goes to post-metadata array
                afterMetadataAndHTML.append(line)

            else:
                inHTML = True

            continue

        elif line.startswith("```"):
            inCode = not inCode
            # afterMetadataAndHTML.append(line)

        elif line.lstrip()[:1] == "#":
            # Heading has triggered end of metadata and end of HTML
            inMetadata = False
            inHTML = False

        elif inHTML:
            continue

        elif inCode:
            afterMetadataAndHTML.append(line)
            continue

        elif line == "\n":
            # Blank line has triggered end of metadata
            inMetadata = False

        if inMetadata is True:
            # Line goes to metadata array
            metadata_lines.append(line)

        else:
            # Line goes to post-metadata array
            afterMetadataAndHTML.append(line)

    want_numbers_headings = False
    want_numbers_content = False

    md2pptx.globals.processingOptions.setOptionValues("slideTemplateFile", "")
    md2pptx.globals.processingOptions.setOptionValues("tempDir", None)
    md2pptx.globals.processingOptions.setOptionValues("hidden",False)

    ######################################################################################
    #                                                                                    #
    # Set default, presentation and current values for some key options                  #
    #                                                                                    #
    ######################################################################################

    md2pptx.globals.processingOptions.setOptionValuesArray(
        [
            ["pageTitleSize", 30],
            ["pageSubtitleSize", "same"],
            ["pageTitleAlign", "left"],
        ]
    )

    md2pptx.globals.processingOptions.setOptionValues("backgroundImage", "")

    # Text size defaults
    md2pptx.globals.processingOptions.setOptionValuesArray(
        [
            ["baseTextSize", 18],
            ["baseTextDecrement", 2],
        ]
    )

    # Code defaults
    md2pptx.globals.processingOptions.setOptionValuesArray(
        [
            ["codeForeground", "000000"],
            ["codeBackground", "DFFFDF"],
            ["codeColumns", 80],
            ["fixedPitchHeightWidthRatio", 1.2],
        ]
    )

    # Text defaults
    md2pptx.globals.processingOptions.setOptionValuesArray(
        [
            ["italicItalic", True],
            ["italicColour", ("None", "")],
            ["boldBold", True],
            ["boldColour", ("None", "")],
        ]
    )

    # Tables defaults
    md2pptx.globals.processingOptions.setOptionValuesArray(
        [
            ["compactTables", 0],
            ["addTableLines", "no"],
            ["addTableColumnLines", []],
            ["addTableRowLines", []],
            ["addTableLineWidth", 1],
            ["addTableLineCount", 1],
            ["addTableLineColour", "000000"],
            ["tableMargin", Inches(0.2)],
            ["spanCells", "yes"],
            ["tableHeadingSize", 0],
            ["tableShadow", False],
        ]
    )

    # Cards defaults
    md2pptx.globals.processingOptions.setOptionValuesArray(
        [
            ["cardPercent", 80],
            ["cardLayout", "horizontal"],
            ["cardTitleAlign", "c"],
            ["cardTitlePosition", "above"],
            ["cardShape", "rounded"],
            ["horizontalCardGap", 0.25],
            ["verticalCardGap", 0.15],
            ["cardShadow", False],
            ["cardTitleSize", 0],
            ["cardBorderWidth", 0],
            ["cardBorderColour", ("None", "")],
            ["cardTitleColour", ("None", "")],
            ["cardColour", [("None", "")]],
            ["cardTitleBackground", [("None", "")]],
            ["cardDividerColour", ("RGB", "#000000")],
            ["cardGraphicSize", 0],
            ["cardGraphicPosition", "before"],
            ["cardGraphicPadding", 0.1],
        ]
    )


    md2pptx.globals.processingOptions.setOptionValues("contentSplit", [1, 1, 1, 1, 1, 1, 1, 1, 1, 1])

    md2pptx.globals.processingOptions.setOptionValues("contentSplitDirection", "vertical")

    # Number of spaces a single level of indentation is
    md2pptx.globals.processingOptions.setOptionValues("indentSpaces", 2)

    # Whether titles are adjusted or not
    md2pptx.globals.processingOptions.setOptionValues("adjustTitles", True)

    # Title and section defaults
    md2pptx.globals.processingOptions.setOptionValuesArray(
        [
            ["sectionTitleSize", 40],
            ["sectionSubtitleSize", 28],
            ["presTitleSize", 40],
            ["presSubtitleSize", 28],
            ["sectionsExpand", False],
        ]
    )

    md2pptx.globals.processingOptions.setOptionValues("monoFont", "Courier")

    topHeadingLevel = 1
    titleLevel = topHeadingLevel
    sectionLevel = titleLevel + 1
    contentLevel = sectionLevel + 1
    cardLevel = contentLevel + 1

    # Abstracts defaults
    abstractTitle = md2pptx.globals.processingOptions.setOptionValues("abstractTitle", "")

    # Tasks defaults
    md2pptx.globals.processingOptions.setOptionValuesArray([["taskSlides", "all"], ["tasksPerPage", 20]])

    # Glossary defaults
    md2pptx.globals.processingOptions.setOptionValuesArray(
        [
            ["glossaryTitle", "Title"],
            ["glossaryTerm", "Term"],
            ["glossaryMeaning", "Meaning"],
            ["glossaryMeaningWidth", 5],
            ["glossaryTermsPerPage", 20],
        ]
    )

    # Footnotes defaults
    md2pptx.globals.processingOptions.setOptionValuesArray(
        [["footnotesTitle", "Footnotes"], ["footnotesPerPage", 20]]
    )

    # Table Of Contents defaults
    md2pptx.globals.processingOptions.setOptionValuesArray(
        [
            ["tocTitle", "Topics"],
            ["tocStyle", ""],
            ["tocItemHeight", 0],
            ["tocItemColour", ""],
            ["tocRowGap", 0.75],
            ["tocFontSize", 0],
            ["tocLinks", False],
            ["sectionArrows", False],
            ["sectionArrowsColour", ""],
        ]
    )

    md2pptx.globals.processingOptions.setOptionValues("marginBase", Inches(0.2))

    md2pptx.globals.processingOptions.setOptionValues("transition", "none")

    md2pptx.globals.processingOptions.setOptionValues("deleteFirstSlide", False)

    md2pptx.globals.processingOptions.setOptionValuesArray(
        [
            ["leftFooterText", ""],
            ["middleFooterText", ""],
            ["rightFooterText", ""],
            ["sectionFooters", "no"],
            ["liveFooters", "no"],
            ["footerFontSize", ""],
        ]
    )

    TOCEntries = []


    metadata = []


    # Space to leave at bottom if numbers
    numbersHeight = Inches(0.4)

    # If no numbers leave all the above height anyway
    md2pptx.globals.processingOptions.setOptionValuesArray(
        [
            ["numbersHeight", numbersHeight],
            ["numbersContentMargin", numbersHeight],
            ["numbersHeadingsMargin", numbersHeight],
            ["numbersFontSize", ""],
        ]
    )

    # Graphics options
    md2pptx.globals.processingOptions.setOptionValuesArray(
        [
            ["exportGraphics", False],
        ]
    )

    # Funnel options
    md2pptx.globals.processingOptions.setOptionValuesArray(
        [
            [
                "funnelColours",
                [
                    ("Theme", MSO_THEME_COLOR.ACCENT_1),
                    ("Theme", MSO_THEME_COLOR.ACCENT_2),
                    ("Theme", MSO_THEME_COLOR.ACCENT_3),
                    ("Theme", MSO_THEME_COLOR.ACCENT_4),
                    ("Theme", MSO_THEME_COLOR.ACCENT_5),
                    ("Theme", MSO_THEME_COLOR.ACCENT_6),
                ],
            ],
            ["funnelBorderColour", ("None", "")],
            ["funnelTitleColour", ("None", "")],
            ["funnelTextColour", ("None", "")],
            ["funnelLabelsPercent", 10],
            ["funnelLabelPosition", "before"],
            ["funnelWidest", "left"],
        ]
    )

    # "on" exit files initialisation
    md2pptx.globals.processingOptions.setOptionValuesArray(
        [
            ["onPresentationInitialisation", ""],
            ["onPresentationBeforeSave", ""],
            ["onPresentationAfterSave", ""],
        ]
    )

    ######################################################################################
    #                                                                                    #
    #  Prime for footnotes                                                               #
    #                                                                                    #
    ######################################################################################

    # List of footnote definitions. Each is a (ref, text) pair.
    # Also array of names - for quick searching
    footnoteDefinitions = []
    footnoteReferences = []

    maxBlocks = 10

    ######################################################################################
    #                                                                                    #
    #  Parse metadata and report on the items found, setting options                     #
    #                                                                                    #
    ######################################################################################

    if len(metadata_lines) > 0:
        print("")
        print("Metadata:")
        print("=========")
        print("")
        print("Name".ljust(40) + " " + "Value")
        print("----".ljust(40) + " " + "-----")


    for line in metadata_lines:
        matchInfo = metadataRegex.match(line)

        if matchInfo is None:
            print("Ignoring invalid metadata line: " + line)
            continue

        name = matchInfo.group(1).strip()
        value = matchInfo.group(2).strip()
        metadata.append([name, value])

        # Print name as it was typed
        print(name.ljust(40) + " " + value)

        # Lower case name for checking
        name = name.lower()

        if name == "numbers":
            numbersHeight = md2pptx.globals.processingOptions.getCurrentOption("numbersHeight")
            if value.lower() == "yes":
                # Want slide numbers everywhere
                want_numbers_headings = True
                md2pptx.globals.processingOptions.setOptionValues("numbersHeadingsMargin", numbersHeight)

                want_numbers_content = True
                md2pptx.globals.processingOptions.setOptionValues("numbersContentMargin", numbersHeight)
            elif value.lower() == "content":
                # Want slide numbers on content slides but not headings & sections
                want_numbers_headings = False
                want_numbers_content = True
                md2pptx.globals.processingOptions.setOptionValues("numbersContentMargin", numbersHeight)
            else:
                # Don't want slide numbers - but they could still be added by slide master
                # (Can code any other value, including 'no' or omit this metadata type)
                want_numbers_headings = False
                want_numbers_content = False

        elif name == "numbersfontsize":
            md2pptx.globals.processingOptions.setOptionValues(name, float(value))

        elif name == "pagetitlesize":
            md2pptx.globals.processingOptions.setOptionValues(name, float(value))

        elif name == "pagetitlealign":
            if value in ["left", "right", "center", "centre", "l", "r", "c"]:
                if value[:1] == "l":
                    md2pptx.globals.processingOptions.setOptionValues(name, "left")
                elif value[:1] == "r":
                    md2pptx.globals.processingOptions.setOptionValues(name, "right")
                elif value[:1] == "c":
                    md2pptx.globals.processingOptions.setOptionValues(name, "center")
            else:
                print(
                    f'PageTitleAlign value \'{value}\' unsupported. "left", "right", "centre", or "center" required.'
                )

        elif name == "pagesubtitlesize":
            if value == "same":
                md2pptx.globals.processingOptions.setOptionValues(name, value)
            else:
                md2pptx.globals.processingOptions.setOptionValues(name, float(value))

        elif name == "basetextsize":
            md2pptx.globals.processingOptions.setOptionValues(name, float(value))

        elif name == "basetextdecrement":
            md2pptx.globals.processingOptions.setOptionValues(name, float(value))

        elif name == "backgroundimage":
            md2pptx.globals.processingOptions.setOptionValues(name, value)

        elif name in [
            "onpresentationinitialisation",
            "onpresentationbeforesave",
            "onpresentationaftersave",
        ]:
            md2pptx.globals.processingOptions.setOptionValues(name, value)
            print(value, "<===")

        elif name in [
            "sectiontitlesize",
            "sectionsubtitlesize",
            "prestitlesize",
            "pressubtitlesize",
        ]:
            md2pptx.globals.processingOptions.setOptionValues(name, float(value))

        elif name == "deletefirstslide":
            if value.lower() == "yes":
                md2pptx.globals.processingOptions.setOptionValues(name, True)
            else:
                md2pptx.globals.processingOptions.setOptionValues(name, False)

        elif name == "hidden":
            if value.lower() == "yes":
                md2pptx.globals.processingOptions.setOptionValues(name, True)
            else:
                md2pptx.globals.processingOptions.setOptionValues(name, False)

        elif name == "sectionsexpand":
            if value.lower() == "yes":
                md2pptx.globals.processingOptions.setOptionValues(name, True)
            else:
                md2pptx.globals.processingOptions.setOptionValues(name, False)

        elif (name == "template") | (name == "master"):
            if value == "Martin Master.pptx":
                slideTemplateFile = "Martin Template.pptx"
            else:
                slideTemplateFile = value
            md2pptx.globals.processingOptions.setOptionValues("slideTemplateFile", slideTemplateFile)

        elif name == "monofont":
            md2pptx.globals.processingOptions.setOptionValues(name, value)

        elif name == "marginbase":
            md2pptx.globals.processingOptions.setOptionValues(name, Inches(float(value)))

        elif name == "tablemargin":
            md2pptx.globals.processingOptions.setOptionValues(name, Inches(float(value)))

        elif name == "tocstyle":
            if value in ["chevron", "circle", "plain"]:
                md2pptx.globals.processingOptions.setOptionValues(name, value)
            else:
                print(
                    f'TOCStyle value \'{value}\' unsupported. "chevron" or "circle" required.'
                )

        elif name == "toctitle":
            md2pptx.globals.processingOptions.setOptionValues(name, value)

        elif name == "tocitemheight":
            md2pptx.globals.processingOptions.setOptionValues(name, float(value))

        elif (name == "tocitemcolour") | (name == "tocitemcolor"):
            md2pptx.globals.processingOptions.setOptionValues("tocItemColour", value)

        elif name == "toclinks":
            if value.lower() == "yes":
                md2pptx.globals.processingOptions.setOptionValues(name, True)

        elif name == "sectionarrows":
            if value.lower() == "yes":
                md2pptx.globals.processingOptions.setOptionValues(name, True)

        elif (name == "sectionarrowscolour") | (name == "sectionarrowscolor"):
            md2pptx.globals.processingOptions.setOptionValues("sectionArrowsColour", value)

        elif name == "tocrowgap":
            md2pptx.globals.processingOptions.setOptionValues(name, float(value))

        elif name == "tocfontsize":
            md2pptx.globals.processingOptions.setOptionValues(name, float(value))

        elif name == "compacttables":
            md2pptx.globals.processingOptions.setOptionValues(name, float(value))

        elif name == "tableheadingsize":
            md2pptx.globals.processingOptions.setOptionValues(name, float(value))

        elif name == "tableshadow":
            if value.lower() == "yes":
                md2pptx.globals.processingOptions.setOptionValues(name, True)

        elif name == "abstracttitle":
            md2pptx.globals.processingOptions.setOptionValues(name, value)

        elif name in [
            "leftfootertext",
            "middlefootertext",
            "rightfootertext",
            "sectionfooters",
            "livefooters",
        ]:
            md2pptx.globals.processingOptions.setOptionValues(name, value)

        elif name == "footerfontsize":
            md2pptx.globals.processingOptions.setOptionValues(name, float(value))

        elif name == "boldbold":
            if value.lower() == "no":
                md2pptx.globals.processingOptions.setOptionValues("boldBold", False)

        elif (name == "boldcolour") | (name == "boldcolor"):
            md2pptx.globals.processingOptions.setOptionValues("boldColour", (parseColour(value.strip())))

        elif name == "italicitalic":
            if value == "no":
                md2pptx.globals.processingOptions.setOptionValues("italicItalic", False)

        elif (name == "italiccolour") | (name == "italiccolor"):
            md2pptx.globals.processingOptions.setOptionValues("italicColour", (parseColour(value.strip())))

        elif name in ["cardcolour", "cardcolor", "cardcolours", "cardcolors"]:
            valueArray2 = [parseColour(c.strip()) for c in value.split(",")]

            md2pptx.globals.processingOptions.setOptionValues("cardColour", valueArray2)

        elif name in ["cardtitlebackground", "cardtitlebackgrounds"]:
            valueArray2 = [parseColour(c.strip()) for c in value.split(",")]

            md2pptx.globals.processingOptions.setOptionValues("cardTitleBackground", valueArray2)

        elif (name == "cardbordercolour") | (name == "cardbordercolor"):
            md2pptx.globals.processingOptions.setOptionValues(
                "cardBorderColour", parseColour(value.strip())
            )

        elif (name == "cardtitlecolour") | (name == "cardtitlecolor"):
            md2pptx.globals.processingOptions.setOptionValues("cardTitleColour", parseColour(value.strip()))

        elif (name == "carddividercolour") | (name == "carddividercolor"):
            md2pptx.globals.processingOptions.setOptionValues(
                "cardDividerColour", parseColour(value.strip())
            )

        elif name == "cardborderwidth":
            md2pptx.globals.processingOptions.setOptionValues(name, float(value))

        elif name == "cardtitlesize":
            md2pptx.globals.processingOptions.setOptionValues(name, float(value))

        elif name == "cardshadow":
            if value.lower() == "yes":
                md2pptx.globals.processingOptions.setOptionValues(name, True)

        elif name in ["cardpercent", "cardgraphicsize", "cardgraphicpadding"]:
            md2pptx.globals.processingOptions.setOptionValues(name, float(value))

        elif name == "cardlayout":
            if value in ["horizontal", "vertical"]:
                md2pptx.globals.processingOptions.setOptionValues(name, value)
            else:
                print(
                    f'CardLayout value \'{value}\' unsupported. "horizontal" or "vertical" required.'
                )

        elif name == "cardshape":
            if value in ["squared", "rounded", "line"]:
                md2pptx.globals.processingOptions.setOptionValues(name, value)
            else:
                print(
                    f'CardShape value \'{value}\' unsupported. "squared", "rounded", or "line" required.'
                )

        elif name == "cardtitleposition":
            if value in ["above", "inside", "before", "after"]:
                md2pptx.globals.processingOptions.setOptionValues(name, value)
            else:
                print(
                    f'CardTitlePosition value \'{value}\' unsupported. "inside", "above", "before", or "after" required.'
                )

        elif name == "cardGraphicPosition":
            if value in ["before", "after"]:
                md2pptx.globals.processingOptions.setOptionValues(name, value)
            else:
                print(
                    f'cardGraphicPosition value \'{value}\' unsupported. "before",or "after" required.'
                )

        elif name == "cardtitlealign":
            val1l = value[:1].lower()
            if val1l in ["l", "r", "c"]:
                md2pptx.globals.processingOptions.setOptionValues(name, val1l)
            else:
                print(f"CardAlign value '{value}' unsupported.")

        elif name in ["horizontalcardgap", "verticalcardgap"]:
            md2pptx.globals.processingOptions.setOptionValues(name, float(value))
            print(float(value))

        elif name == "contentsplit":
            splitValue = value.split()
            cs = []
            for v in splitValue:
                cs.append(int(v))

            # Extend to maximum allowed
            needMore = maxBlocks - len(cs)
            for _ in range(needMore):
                cs.append(1)

            md2pptx.globals.processingOptions.setOptionValues("contentSplit", cs)

        elif name in ["contentsplitdirection", "contentsplitdirn"]:
            if value in [
                "vertical",
                "horizontal",
                "v",
                "h",
                "default",
                "pres",
                "pop",
                "prev",
            ]:
                if value in ["vertical", "horizontal"]:
                    adjustedValue = value
                elif value == "v":
                    adjustedValue = "vertical"
                elif value == "h":
                    adjustedValue = "horizontal"
                else:
                    adjustedValue = value

                md2pptx.globals.processingOptions.setOptionValues("contentSplitDirection", adjustedValue)

            else:
                print(
                    f'{name} value \'{value}\' unsupported. "vertical" or "horizontal" required.'
                )

        elif name == "taskslides":
            md2pptx.globals.processingOptions.setOptionValues(name, value)

        elif name == "tasksperpage":
            md2pptx.globals.processingOptions.setOptionValues(name, int(value))

        elif name in [
            "titleslidelayout",
            "sectionslidelayout",
            "contentslidelayout",
            "titleonlylayout",
            "blanklayout",
        ]:
            md2pptx.globals.processingOptions.setOptionValues(name, int(value))

        elif name == "numbersheight":
            numbersHeight = Inches(float(value))

            # If no numbers leave all the above height anyway
            md2pptx.globals.processingOptions.setOptionValues("numbersHeight", numbersHeight)
            md2pptx.globals.processingOptions.setOptionValues("numbersContentMargin", numbersHeight)
            md2pptx.globals.processingOptions.setOptionValues("numbersHeadingsMargin", numbersHeight)

        elif name == "glossarytitle":
            md2pptx.globals.processingOptions.setOptionValues(name, value)

        elif name == "glossaryterm":
            md2pptx.globals.processingOptions.setOptionValues(name, value)

        elif name == "glossarymeaning":
            md2pptx.globals.processingOptions.setOptionValues(name, value)

        elif name == "glossarymeaningwidth":
            md2pptx.globals.processingOptions.setOptionValues(name, int(value))

        elif name == "glossarytermsperpage":
            md2pptx.globals.processingOptions.setOptionValues(name, int(value))

        elif name == "footnotesperpage":
            md2pptx.globals.processingOptions.setOptionValues(name, int(value))

        elif name == "footnotestitle":
            md2pptx.globals.processingOptions.setOptionValues(name, value)

        # Following relate to styling and don't use Processing Options class

        elif name.startswith("style.bgcolor."):
            spanClass = name[14:]
            if value.startswith("#"):
                value2 = value
            else:
                value2 = "#" + value
            (check, colour) = parseRGB(value2)
            if check:
                # Valid RGB hex value
                md2pptx.globals.bgcolors[spanClass] = colour
            else:
                print(f"Invalid value for {name} - {value}. Ignoring.")

        elif name.startswith("style.fgcolor."):
            spanClass = name[14:]
            if value.startswith("#"):
                value2 = value
            else:
                value2 = "#" + value
            (check, colour) = parseRGB(value2)
            if check:
                # Valid RGB hex value
                md2pptx.globals.fgcolors[spanClass] = colour
            else:
                print(f"Invalid value for {name} - {value}. Ignoring.")

        elif name.startswith("style.emphasis."):
            spanClass = name[15:]
            md2pptx.globals.emphases[spanClass] = value

        elif name.startswith("style.fontsize."):
            spanClass = name[15:]

            # Assumed "px" on the end
            value2 = value[:-2]
            md2pptx.globals.fontsizes[spanClass] = value2

        elif name in ["codeforeground", "codebackground"]:
            md2pptx.globals.processingOptions.setOptionValues(name, value)

        elif name == "fpratio":
            md2pptx.globals.processingOptions.setOptionValues("fixedPitchHeightWidthRatio", float(value))

        elif name == "codecolumns":
            md2pptx.globals.processingOptions.setOptionValues(name, int(value))

        elif name == "topheadinglevel":
            titleLevel = int(value)
            sectionLevel = titleLevel + 1
            contentLevel = sectionLevel + 1
            cardLevel = contentLevel + 1

        elif name == "indentspaces":
            md2pptx.globals.processingOptions.setOptionValues(name, int(value))

        elif name == "adjusttitles":
            if value == "no":
                md2pptx.globals.processingOptions.setOptionValues(name, False)

        elif name in ["addtablelines", "addtablelinecolour"]:
            md2pptx.globals.processingOptions.setOptionValues(name, value)

        elif name in ["addtablecolumnlines", "addtablerowlines"]:
            md2pptx.globals.processingOptions.setOptionValues(name, sortedNumericList(value))

        elif name in ["addtablelinecount", "addtablelinewidth"]:
            md2pptx.globals.processingOptions.setOptionValues(name, int(value))

        elif name == "spancells":
            md2pptx.globals.processingOptions.setOptionValues(name, value)

        elif name == "hidemetadata":
            if value == "style":
                md2pptx.globals.processingOptions.hideMetadataStyle = True

        elif name == "exportgraphics":
            if value.lower() == "yes":
                md2pptx.globals.processingOptions.setOptionValues(name, True)

        elif name == "tempdir":
            md2pptx.globals.processingOptions.setOptionValues(name, os.path.expanduser(value))

        elif name == "transition":
            if value.lower() in [
                "none",
                "ripple",
                "reveal",
                "honeycomb",
                "shred",
                "wipe",
                "push",
                "vortex",
                "split",
                "fracture",
            ]:
                # Valid transition name
                md2pptx.globals.processingOptions.setOptionValues(name, value.lower())
            else:
                print(f"Invalid value for {name} - {value}. Ignoring.")

        elif (name == "funnelcolours") | (name == "funnelcolors"):
            valueArray2 = [parseColour(c.strip()) for c in value.split(",")]

            md2pptx.globals.processingOptions.setOptionValues("funnelColours", valueArray2)

        elif (name == "funnelbordercolour") | (name == "funnelbordercolor"):
            md2pptx.globals.processingOptions.setOptionValues(
                "funnelBorderColour", parseColour(value.strip())
            )

        elif (name == "funneltitlecolour") | (name == "funneltitlecolor"):
            md2pptx.globals.processingOptions.setOptionValues(
                "funnelTitleColour", parseColour(value.strip())
            )

        elif (name == "funneltextcolour") | (name == "funneltextcolor"):
            md2pptx.globals.processingOptions.setOptionValues(
                "funnelTextColour", parseColour(value.strip())
            )

        elif name == "funnellabelspercent":
            md2pptx.globals.processingOptions.setOptionValues(name, float(value))

        elif name == "funnellabelposition":
            if value in ["before", "after"]:
                md2pptx.globals.processingOptions.setOptionValues(name, value)
            else:
                print(
                    f'funnelLabelPosition value \'{value}\' unsupported. "before" or "after" required.'
                )

        elif name == "funnelwidest":
            if value in ["left", "right", "pipe", "hpipe", "vpipe", "top", "bottom"]:
                md2pptx.globals.processingOptions.setOptionValues(name, value)
            else:
                print(
                    f'funnelLabelPosition value \'{value}\' unsupported. "left", "right", "pipe", or "hpipe" required.'
                )

    slideTemplateFile = md2pptx.globals.processingOptions.getCurrentOption("slideTemplateFile")
    if slideTemplateFile != "":
        originalSlideTemplateFile = slideTemplateFile
        if Path(os.path.expanduser(slideTemplateFile)).exists():
            # We can successfully expand the path to pick up the file
            slideTemplateFile = os.path.expanduser(slideTemplateFile)
        else:
            # Slide template file is not present if we expand path
            script_path = os.path.dirname(__file__)
            slideTemplateFile = script_path + os.sep + slideTemplateFile
            if not Path(slideTemplateFile).exists():
                print(
                    f"\nTemplate file {originalSlideTemplateFile} does not exist. Terminating."
                )
                sys.exit()

        print(f"\nUsing {slideTemplateFile} as base for presentation")

    if slideTemplateFile == "":
        # Use default slide deck that comes with python-pptx as base
        prs = Presentation()
        print("\nNo slide to document metadata on. Continuing without it.")

        templateSlideCount = 0
    else:
        # Use user-specified presentation as base
        prs = Presentation(slideTemplateFile)

        # If there is a slide to use fill it with metadata
        templateSlideCount = len(prs.slides)
        if templateSlideCount > 0:
            print("\nWriting processing summary slide with metadata on it.")
        else:
            print("\nNo slide to document metadata on. Continuing without it.")

        # Prime template slides with slideInfo as None
        for slide in prs.slides:
            slide.slideInfo = None

    # Maybe call an exit as the presentation is initialised
    onPresentationInitialisation = md2pptx.globals.processingOptions.getCurrentOption("onPresentationInitialisation")
    if onPresentationInitialisation != "":
        exec(open(onPresentationInitialisation).read())

    # Following might be used in slide footers
    prs.lastSectionTitle = ""
    prs.lastSectionSlide = None
    prs.lastPresTitle = ""
    prs.lastPresSubtitle = ""

    print("")
    print("Slides:")
    print("=======")
    print("")

    inBlock = False
    inList = False
    inTable = False
    inCard = False
    inTitle = False

    blockType = ""
    slideTitle = ""
    slideSubtitle = ""
    bullets = []
    tableRows = []
    cards = []
    code = []
    inCode = False
    inHTMLCode = False
    inFencedCode = False
    notes_text = ""
    slide = None
    tasks = []
    sequence = []

    slideHrefs = {}

    # Each of these is a picture, then a href, then a tooltip - as a tuple
    pictureInfos = []

    # Pass 2: Concatenate lines with continuations
    previousLine = "\n"
    linesAfterConcatenation = []

    for line in afterMetadataAndHTML:
        if startswithOneOf(line, ["<pre>", "<code>"]):
            # These are around code lines
            linesAfterConcatenation.append(line)
            inHTMLCode = True

        elif startswithOneOf(line, ["</pre>", "</code>"]):
            # These are around code lines
            linesAfterConcatenation.append(line)
            inHTMLCode = False

        elif line.startswith("```"):
            linesAfterConcatenation.append(line)
            inCode = not inCode

        elif line == "\n":
            # This is a blank line so copy it over
            linesAfterConcatenation.append(line)

        elif previousLine == "\n":
            # Previous line was blank so copy this one over
            linesAfterConcatenation.append(line)

        elif line.startswith("<!-- md2pptx: "):
            # This is a dynamic metadata line so keep it separate
            linesAfterConcatenation.append(line)

        elif line.startswith("<a id="):
            # This is an anchor line so keep it separate
            linesAfterConcatenation.append(line)

        elif startswithOneOf(line, ["<video ", "<audio "]):
            # This is a video / audio element line
            linesAfterConcatenation.append(line)

        elif line.lstrip() == "":
            # This is an empty line
            linesAfterConcatenation.append(line)

        else:
            # Previous line was not blank and nor is this one so consider concatenation
            if line.lstrip()[0] not in r"*#\|0123456789!":
                if (
                    (previousLine[0:2] != "# ")
                    & (previousLine[0:3] != "## ")
                    & (previousLine[0:4] != "    ")
                    & (previousLine[0] != "|")
                    & (inCode is False)
                    & (inHTMLCode is False)
                ):
                    # Previous line was not Heading Level 1 or 2 and we're not in code so concatenate
                    linesAfterConcatenation[-1] = (
                        linesAfterConcatenation[-1].rstrip() + " " + line.lstrip()
                    )
                else:
                    linesAfterConcatenation.append(line)

            else:
                linesAfterConcatenation.append(line)

        # Store previous line to see if it was H1 or blank
        previousLine = line

    # Pass 3: Get footnote definitions
    metadataLinenumber = 0
    for line in linesAfterConcatenation:
        line = line.rstrip()
        if m := footnoteDefinitionRegex.match(line):
            fnRef = m.group(1).strip()
            fnText = m.group(2).strip()
            footnoteDefinitions.append([fnRef, fnText])
            footnoteReferences.append(fnRef)

            linesAfterConcatenation[metadataLinenumber] = "<ignoreme>"
        metadataLinenumber += 1

    # Pass 4: Extract any indirect reference anchors
    metadataLinenumber = 0
    indirectAnchors = []
    for line in linesAfterConcatenation:
        line = line.rstrip()
        if m := indirectReferenceAnchorRegex.match(line):
            anchorName = m.group(1).strip()
            anchorURL = m.group(2).strip()

            indirectAnchors.append([anchorName, anchorURL])

            linesAfterConcatenation[metadataLinenumber] = "<ignoreme>"
        metadataLinenumber += 1

    lastTableLine = -1
    tableCaptions = []

    # Pass 5: Main pass over the input file, now that footnote
    # references have been gathered
    for lineNumber, line in enumerate(linesAfterConcatenation):
        # Remove trailing white space
        line = line.rstrip()

        # Convert tabs to spaces
        line = line.replace("\t", " " * md2pptx.globals.processingOptions.getCurrentOption("indentSpaces"))

        if line == "<ignoreme>":
            # Line was taken care of in the previous pass
            continue

        if startswithOneOf(line, ["<pre>", "<code>"]):
            code.append([])
            inCode = True
            inHTMLCode = True
            inTable = False
            inTitle = False

        if startswithOneOf(line, ["</pre>", "</code>"]):
            inCode = False
            inHTMLCode = False
            inTitle = False

        if line.startswith("```"):
            inCode = not inCode
            if inCode:
                # Just entered code
                code.append([])
                blockType = "code"
                inTable = False
            else:
                # Just exited code - but add closing line
                code[-1].append(line)
            inFencedCode = not inFencedCode
            inTitle = False

        if inCode or inHTMLCode or inFencedCode:
            if len(code) == 0:
                code.append([])
            code[-1].append(line)

            # If first line of code then mark the current sequence entry as "code"
            if len(code[-1]) == 1:
                sequence.append("code")

            inTitle = False
        if (
            (line == "")
            & (inCode is True)
            & (inHTMLCode is False)
            & (inFencedCode is False)
        ):
            inCode = False
            inTitle = False

        if (line.startswith("    ")) & (inList is False):
            # Only list items and code can be indented by 4 characters
            if inCode is False:
                code.append([])
                blockType = "code"
                code[-1].append(line[4:])
                inCode = True
            inTitle = False

        # Rewrite horizontal rule as a heading 3 with non-breaking space
        if startswithOneOf(line, ["<hr/>", "---", "***", "___"]):
            line = "### &nbsp;"
            inTitle = True

        # Taskpaper task
        if line[:1] == "-":
            # Get start of attributes
            attributesStart = line.find("@")

            # Get text up to attributes
            if attributesStart == -1:
                text = line[2:]
            else:
                text = line[2 : attributesStart - 1]

            # Attempt to extract @due information
            startDue = line.find("@due(")
            if startDue > -1:
                startDue += 5
                endDue = line.find(")", startDue)
                if endDue > -1:
                    dueDate = line[startDue:endDue]
            else:
                dueDate = ""

            # Attempt to extract @tags information
            startTags = line.find("@tags(")
            if startTags > -1:
                startTags += 6
                endTags = line.find(")", startTags)
                if endTags > -1:
                    tags = line[startTags:endTags]
            else:
                tags = ""

            # Attempt to extract @done information
            startDone = line.find("@done(")
            if startDone > -1:
                startDone += 6
                endDone = line.find(")", startDone)
                if endDone > -1:
                    done = line[startDone:endDone]
            else:
                done = ""

            tasks.append([slideNumber + 1, text, dueDate, tags, done])
            inTitle = False

        elif line.startswith("<a id="):
            # Anchor on whatever slide we're on
            if hrefMatch := anchorRegex.match(line):
                href = hrefMatch.group(1)
                if (href != "") & (href in slideHrefs.keys()):
                    print(f"Heading Reference redefined: '{href}' for slide {slideNumber}")
            inTitle = False

        elif line.startswith("<!-- md2pptx: "):
            # Dynamic metadata line
            inTitle = False
            if DMMatch := dynamicMetadataRegex.match(line):
                metadataKey = DMMatch.group(1).lower().strip()
                metadataValue = DMMatch.group(2).lower()
                if (metadataKey != "") & (metadataValue != ""):
                    # Valid metadata pair so parse key / value - and apply if valid

                    if (metadataKey == "pagesubtitlesize") & (metadataValue == "same"):
                        # Cards' layout - horizontal or vertical
                        md2pptx.globals.processingOptions.dynamicallySetOption(
                            metadataKey, metadataValue, ""
                        )

                    # Floating point values where metadata key matches directly
                    elif metadataKey in [
                        # Font size for tables
                        "compacttables",
                        # Heading font size for tables
                        "tableheadingsize",
                        # Page Title Font size
                        "pagetitlesize",
                        # Page Subtitle Font size
                        "pagesubtitlesize",
                        # Cards' vertical percent of content area
                        "cardpercent",
                        # Card graphic's vertical or horizontal space
                        "cardgraphicsize",
                        # Base text size
                        "basetextsize",
                        # Base text decrement
                        "basetextdecrement",
                        # Horizontal card gap
                        "horizontalcardgap",
                        # Vertical card gap
                        "verticalcardgap",
                        # Funnel label percentage
                        "funnellabelspercent",
                    ]:
                        md2pptx.globals.processingOptions.dynamicallySetOption(
                            # Use the metadata key directly
                            metadataKey,
                            metadataValue.lower(),
                            "float",
                        )

                    elif metadataKey == "cardlayout":
                        # Cards' layout - horizontal or vertical
                        md2pptx.globals.processingOptions.dynamicallySetOption(
                            metadataKey, metadataValue, ""
                        )
                    elif metadataKey == "hidden":
                        # Slides hidden in slideshow
                        if metadataValue == "yes":
                            metadataValue = True
                        elif metadataValue == "no":
                            metadataValue = False

                        md2pptx.globals.processingOptions.dynamicallySetOption(
                            metadataKey, metadataValue, ""
                        )

                    elif metadataKey == "numbersheight":
                        # vertical space to reserve for slide number
                        if metadataValue == "default":
                            numbersheight = md2pptx.globals.processingOptions.getDefaultOption(
                                "numbersHeight"
                            )
                        elif metadataValue == "pres":
                            numbersheight = md2pptx.globals.processingOptions.getPresentationOption(
                                "numbersHeight"
                            )
                        else:
                            numbersheight = int(Inches(float(metadataValue)))

                        md2pptx.globals.processingOptions.dynamicallySetOption(
                            metadataKey, numbersheight, "int"
                        )

                    elif metadataKey == "marginbase":
                        #  space to reserve as a margin
                        if metadataValue == "default":
                            marginbase = md2pptx.globals.processingOptions.getDefaultOption("marginbase")
                        elif metadataValue == "pres":
                            marginbase = md2pptx.globals.processingOptions.getPresentationOption(
                                "marginbase"
                            )
                        else:
                            marginbase = int(Inches(float(metadataValue)))

                        md2pptx.globals.processingOptions.dynamicallySetOption(
                            metadataKey, marginbase, "int"
                        )

                    elif metadataKey == "tablemargin":
                        #  space to reserve as a table margin
                        if metadataValue == "default":
                            tablemargin = md2pptx.globals.processingOptions.getDefaultOption("tablemargin")
                        elif metadataValue == "pres":
                            tablemargin = md2pptx.globals.processingOptions.getPresentationOption(
                                "tablemargin"
                            )
                        else:
                            tablemargin = int(Inches(float(metadataValue)))

                        md2pptx.globals.processingOptions.dynamicallySetOption(
                            metadataKey, tablemargin, "int"
                        )

                    elif metadataKey == "spancells":
                        #  whether table cells can span more than one column
                        if metadataValue == "default":
                            spancells = md2pptx.globals.processingOptions.getDefaultOption("spancells")
                        elif metadataValue == "pres":
                            spancells = md2pptx.globals.processingOptions.getPresentationOption("spancells")
                        else:
                            spancells = metadataValue

                        md2pptx.globals.processingOptions.dynamicallySetOption(metadataKey, spancells, "")

                    # Note: Actual value handling prevents using dynamicallySetOption
                    elif metadataKey == "cardtitlealign":
                        # Card title alignment
                        md2pptx.globals.processingOptions.dynamicallySetOption(
                            metadataKey, metadataValue.lower(), ""
                        )

                    elif metadataKey == "cardtitleposition":
                        # Card title position - above or inside
                        md2pptx.globals.processingOptions.dynamicallySetOption(
                            metadataKey, metadataValue, ""
                        )

                    elif metadataKey == "cardGraphicPosition":
                        # Card graphic position - before or after
                        md2pptx.globals.processingOptions.dynamicallySetOption(
                            metadataKey, metadataValue, ""
                        )

                    elif metadataKey == "cardshape":
                        # Card shape
                        if metadataValue in ["squared", "rounded", "line"]:
                            md2pptx.globals.processingOptions.dynamicallySetOption(
                                metadataKey, metadataValue, ""
                            )
                        else:
                            print(
                                f'CardShape value \'{metadataValue}\' unsupported. "squared", "rounded", or "line" required.'
                            )

                    elif metadataKey in [
                        "cardcolour",
                        "cardcolor",
                        "cardcolours",
                        "cardcolors",
                    ]:
                        valueArray2 = [
                            parseColour(c.strip()) for c in metadataValue.split(",")
                        ]

                        md2pptx.globals.processingOptions.dynamicallySetOption(
                            "cardColour", valueArray2, ""
                        )

                    elif metadataKey in [
                        "cardtitlebackground",
                        "cardtitlebackgrounds",
                    ]:
                        valueArray2 = [
                            parseColour(c.strip()) for c in metadataValue.split(",")
                        ]

                        md2pptx.globals.processingOptions.dynamicallySetOption(
                            "cardTitleBackground", valueArray2, ""
                        )

                    elif metadataKey in ["funnelcolours", "funnelcolors"]:
                        valueArray2 = [
                            parseColour(c.strip()) for c in metadataValue.split(",")
                        ]

                        md2pptx.globals.processingOptions.dynamicallySetOption(
                            "funnelColours", valueArray2, ""
                        )

                    elif metadataKey in [
                        "funneltitlecolour",
                        "funnelbordercolour",
                        "funneltextcolour",
                    ]:
                        # Funnel single colour options
                        md2pptx.globals.processingOptions.dynamicallySetOption(
                            metadataKey, metadataValue, ""
                        )

                    elif metadataKey == "backgroundimage":
                        # Slide background image
                        md2pptx.globals.processingOptions.dynamicallySetOption(
                            metadataKey, metadataValue, ""
                        )


                    # Note: Actual value handling prevents using dynamicallySetOption
                    elif metadataKey == "contentsplit":
                        # Proportions for each content element on a slide
                        if metadataValue == "default":
                            contentSplit = md2pptx.globals.processingOptions.getDefaultOption(
                                "contentSplit"
                            )
                        elif metadataValue == "pres":
                            contentSplit = md2pptx.globals.processingOptions.getPresentationOption(
                                "contentSplit"
                            )
                        elif metadataValue in ["pop", "prev"]:
                            md2pptx.globals.processingOptions.popCurrentOption("contentSplit")
                            contentSplit = md2pptx.globals.processingOptions.getPresentationOption(
                                "contentSplit"
                            )
                        else:
                            splitValue = metadataValue.split()
                            contentSplit = []
                            for v in splitValue:
                                contentSplit.append(int(v))

                            # Extend to maximum allowed
                            needMore = maxBlocks - len(contentSplit)
                            for _ in range(needMore):
                                contentSplit.append(1)
                        md2pptx.globals.processingOptions.dynamicallySetOption(
                            "contentSplit", contentSplit, ""
                        )

                    elif metadataKey in ["contentsplitdirection", "contentsplitdirn"]:
                        if metadataValue in [
                            "vertical",
                            "horizontal",
                            "v",
                            "h",
                            "pres",
                            "default",
                            "pop",
                            "prev",
                        ]:
                            if metadataValue == "v":
                                adjustedValue = "vertical"
                            elif metadataValue == "h":
                                adjustedValue = "horizontal"
                            else:
                                adjustedValue = metadataValue

                            md2pptx.globals.processingOptions.dynamicallySetOption(
                                "contentSplitDirection",
                                adjustedValue,
                                "",
                            )

                        else:
                            print(
                                f'{metadataKey} value \'{metadataValue}\' unsupported. "vertical" or "horizontal" required.'
                            )

                    elif metadataKey in ["codeforeground", "codebackground"]:
                        md2pptx.globals.processingOptions.dynamicallySetOption(
                            metadataKey,
                            metadataValue,
                            "",
                        )

                    elif metadataKey == "fpratio":
                        # Fixed Pitch font height to width ratio
                        md2pptx.globals.processingOptions.dynamicallySetOption(
                            "fixedPitchHeightWidthRatio",
                            metadataValue,
                            "float",
                        )

                    elif metadataKey == "codecolumns":
                        md2pptx.globals.processingOptions.dynamicallySetOption(
                            metadataKey,
                            metadataValue,
                            "int",
                        )

                    elif metadataKey == "indentspaces":
                        # Spaces representing a single indentation level
                        md2pptx.globals.processingOptions.dynamicallySetOption(
                            metadataKey,
                            metadataValue,
                            "int",
                        )

                    elif metadataKey in ["addtablecolumnlines", "addtablerowlines"]:
                        md2pptx.globals.processingOptions.dynamicallySetOption(
                            metadataKey,
                            metadataValue,
                            "sortednumericlist",
                        )

                    elif metadataKey in ["addtablelinecount", "addtablelinewidth"]:
                        md2pptx.globals.processingOptions.dynamicallySetOption(
                            metadataKey,
                            metadataValue,
                            "int",
                        )

                    elif metadataKey in ["addtablelines", "addtablelinecolour"]:
                        md2pptx.globals.processingOptions.dynamicallySetOption(
                            metadataKey,
                            metadataValue,
                            "",
                        )

                    elif metadataKey == "transition":
                        if metadataValue.lower() in [
                            "none",
                            "ripple",
                            "reveal",
                            "honeycomb",
                            "shred",
                            "wipe",
                            "push",
                            "vortex",
                            "split",
                            "fracture",
                            "pres",
                            "default",
                            "pop",
                            "prev",
                        ]:
                            md2pptx.globals.processingOptions.dynamicallySetOption(
                                metadataKey,
                                metadataValue,
                                "",
                            )
                        else:
                            print(
                                f"Invalid value for {metadataKey}: {metadataValue} in '{line}"
                            )

                    elif metadataKey == "funnellabelposition":
                        if metadataValue.lower() in [
                            "before",
                            "after",
                        ]:
                            md2pptx.globals.processingOptions.dynamicallySetOption(
                                metadataKey,
                                metadataValue,
                                "",
                            )
                        else:
                            print(
                                f"Invalid value for {metadataKey}: {metadataValue} in '{line}"
                            )

                    elif metadataKey == "funnelwidest":
                        if metadataValue.lower() in [
                            "left",
                            "right",
                            "pipe",
                            "hpipe",
                            "vpipe",
                            "top",
                            "bottom",
                        ]:
                            md2pptx.globals.processingOptions.dynamicallySetOption(
                                metadataKey,
                                metadataValue,
                                "",
                            )
                        else:
                            print(
                                f"Invalid value for {metadataKey}: {metadataValue} in '{line}"
                            )

                    else:
                        # Invalid dynamic metadata specification
                        print(f"Invalid dynamic metadata key: '{metadataKey}' in '{line}'")

        elif (line.startswith("#")) & (inCode == False):
            if line[:cardLevel] == "#" * cardLevel:
                # Card on an existing slide
                inCard = True

                # If there is an outstanding href then store it
                if href != "":
                    slideHrefs[href] = slideNumber + 1

                # Get card title text and any href
                cardName, href = parseTitleText(line[cardLevel:])

                # No bullets for the card yet
                cardBullets = []

                # No graphic for the card yet
                cardGraphic = ""

                # Create information about a card_title
                card = Card()
                card.title = cardName
                card.bullets = cardBullets
                card.graphic = cardGraphic
                cards.append(card)

                # Might have octothorpes but isn't a title
                inTitle = False


            else:
                # One of Content, Section, or Title slide
                if inBlock is True:
                    # Create the previous slide
                    slideInfo = SlideInfo(
                        slideTitle,
                        slideSubtitle,
                        blockType,
                        bullets,
                        tableRows,
                        cards,
                        code,
                        sequence,
                    )
                    slideNumber, slide, sequence = createSlide(prs, slideNumber, slideInfo)

                    # Register the previous slide's href - if there was one
                    if href != "":
                        slideHrefs[href] = slideNumber

                if line[:contentLevel] == "#" * contentLevel:
                    # Heading Level 3 - slide
                    thisLevel = contentLevel
                    blockType = "content"

                elif line[:sectionLevel] == "#" * sectionLevel:
                    # Heading Level 2 - section
                    thisLevel = sectionLevel
                    blockType = "section"
                    inTitle = True

                elif line[:titleLevel] == "#" * titleLevel:
                    # Heading Level 1 - slide title
                    thisLevel = titleLevel
                    blockType = "title"
                    inTitle = True
                else:
                    inTitle = False

                # Get slide title text and any href
                slideTitle, href = parseTitleText(line[thisLevel:])

                inBlock = True
                inList = False
                inTable = False
                inCard = False
                slideSubtitle = ""
                bullets = []
                tableRows = []
                cards = []
                code = []

                if (notes_text != "") & (slide != None):
                    createSlideNotes(slide, notes_text)

                notes_text = ""

            # Check whether any heading reference is a duplicate
            if (href != "") & (href in slideHrefs.keys()):
                print(f"Heading Reference redefined: '{href}' for slide {slideNumber}")

        elif match := bulletRegex.match(line):
            # Bulleted list
            bulletLine = match.group(3).lstrip()

            bulletLevel = calculateIndentationLevel(
                match.start(2), md2pptx.globals.processingOptions.getCurrentOption("indentSpaces")
            )

            bulletType = "bulleted"

            if inCard:
                # Bullet is on a card
                cardBullets.append([bulletLevel, bulletLine, bulletType])

                # Update bullet list for the card
                cards[-1].bullets = cardBullets
            else:
                # Bullet is not on a card
                bullets.append([bulletLevel, bulletLine, bulletType])

            if inList is False:
                sequence.append("list")

            inList = True
            inTable = False
            inTitle = False

        elif match := numberRegex.match(line):
            # Numbered list
            bulletLine = match.group(3).lstrip()

            bulletLevel = calculateIndentationLevel(
                match.start(2), md2pptx.globals.processingOptions.getCurrentOption("indentSpaces")
            )

            bulletType = "numbered"

            if inCard:
                # Bullet is on a card
                cardBullets.append([bulletLevel, bulletLine, bulletType])

                # Update bullet list for the card
                cards[-1].bullets = cardBullets
            else:
                # Bullet is not on a card
                bullets.append([bulletLevel, bulletLine, bulletType])

            if inList is False:
                sequence.append("list")

            inList = True
            inTable = False
            inTitle = False

        # Following marshals media into one or two table cells and adds the new row
        # to either a new or existing table. It also says we're in a table
        elif startswithOneOf(line, ["<video ", "<audio ", "[![", "!["]):
            if inCard:
                # Graphic in a card
                if startswithOneOf(line, ["<video ", "<audio ", "[!["]):
                    if clickableGraphicMatch := clickableGraphicRegex.match(line):
                        # Line contains a clickable graphic
                        cards[-1].graphicTitle = clickableGraphicMatch.group(1)

                        (
                            cards[-1].graphic,
                            cards[-1].printableFilename,
                        ) = handleWhateverGraphicType(clickableGraphicMatch.group(2))

                        cards[-1].mediaURL = clickableGraphicMatch.group(3)

                    elif videoRegexMatch := videoRegex.match(line):
                        (
                            _ ,
                            _ ,
                            cards[-1].printableFilename,
                            _ ,
                            _ ,
                            cards[-1].mediaInfo,
                            _ ,
                        ) = parseMedia(videoRegexMatch.group(0), 0)

                    else:
                        audioRegexMatch = audioRegex.match(line)
                        (
                            _ ,
                            _ ,
                            cards[-1].printableFilename,
                            _ ,
                            _ ,
                            cards[-1].mediaInfo,
                            _ ,
                        ) = parseMedia(audioRegexMatch.group(0), 0)

                else:
                    # Cell contains a non-clickable graphic
                    graphicMatch = graphicRegex.match(line)

                    cards[-1].graphicTitle = graphicMatch.group(1)

                    (
                        cards[-1].graphic,
                        cards[-1].printableFilename,
                    ) = handleWhateverGraphicType(graphicMatch.group(2))
            else:
                # There is at least one media item in the line and not in a card
                # - so set up table
                if inTable is False:
                    # Switch to being in a table
                    blockType = "table"
                    sequence.append("table")
                    inTable = True
                    inList = False
                    inCard = False

                    # Create a new empty table
                    # tableRows = []

                # Start the creation of a new row
                tableRow = []

                # Collect the media items
                mediaItems = []
                for m in re.finditer(videoRegex, line):
                    # A video
                    mediaItem = [m.start(0), m.end(0), m.group(0)]
                    mediaItems.append(mediaItem)

                for m in re.finditer(audioRegex, line):
                    # An audio file
                    mediaItem = [m.start(0), m.end(0), m.group(0)]
                    mediaItems.append(mediaItem)

                for m in re.finditer(clickableGraphicRegex, line):
                    # A clickable graphic
                    mediaItem = [m.start(0), m.end(0), m.group(0)]
                    mediaItems.append(mediaItem)

                for m in re.finditer(graphicRegex, line):
                    # A non-clickable graphic
                    start = m.start(0)
                    end = m.end(0)
                    mediaItem = [start, end, m.group(0)]

                    rematch = False
                    for existingMediaItem in mediaItems:
                        if (start >= existingMediaItem[0]) & (end <= existingMediaItem[1]):
                            # This graphic already found as part of a clickable graphic link
                            rematch = True
                            break

                    if rematch is False:
                        # This graphic not already found as part of a clickable graphic link
                        mediaItems.append(mediaItem)

                # compose table row - based on start offset
                tableRow = []
                for mediaItem in sorted(mediaItems):
                    tableRow.append(mediaItem[2])

                # Add table row to saved table rows
                if len(tableRows) == 0:
                    tableRows.append([])
                tableRows[-1].append(tableRow)

                inTitle = False

        elif line[:1] == "|":
            # Table or side-by-side
            lastTableLine = lineNumber

            # As we're in a table we can't have a table caption yet
            tableCaption = ""

            if inTable is False:
                blockType = "table"
                tableRows.append([])
                inTable = True
                sequence.append("table")
                inList = False
                inCard = False

            # Create a table row - but with (maybe empty) junk before and after
            words = line.split("|")
            tableRow = []
            for cell in words:
                tableRow.append(cell)

            # Remove first element
            tableRow.pop(0)

            # Remove last element - if blank
            if cell == "":
                tableRow.pop()

            # Add clean table row to saved table rows
            tableRows[-1].append(tableRow)
            inTitle = False

        elif (
            (line.startswith("[")) & line.endswith("]") & (lineNumber == lastTableLine + 1)
        ):
            tableCaption = line[1:-1]
            tableCaptions.append(tableCaption)
        else:
            # Not in a table
            inTable = False

            if len(tableCaptions) < len(tableRows):
                tableCaptions.append("")

            if not inCode:
                # Must be a slide note line or a subtitle line
                if line == "":
                    inTitle = False
                if inTitle:
                    slideSubtitle = slideSubtitle + "\n" + line
                elif startswithOneOf(line, ["</pre>", "</code>"]) is False:
                    notes_text = notes_text + "\n" + line

    # Ensure there's a blank table caption - for the case the table ended in the final line
    if len(tableCaptions) < len(tableRows):
        tableCaptions.append(tableCaption)

    ######################################################################################
    #                                                                                    #
    # Finish off last slide                                                              #
    #                                                                                    #
    ######################################################################################
    if (inBlock is True) | (inCode is True) | (inTable is True):
        slideInfo = SlideInfo(
            slideTitle, slideSubtitle, blockType, bullets, tableRows, cards, code, sequence
        )
        slideNumber, slide, sequence = createSlide(prs, slideNumber, slideInfo)

        if href != "":
            slideHrefs[href] = slideNumber

        if (notes_text != "") & (slide != None):
            createSlideNotes(slide, notes_text)

        notes_text = ""

    ######################################################################################
    #                                                                                    #
    # Add a footnotes slide - if there were any footnote definitions                     #
    #                                                                                    #
    ######################################################################################
    if len(footnoteDefinitions) > 0:
        slideNumber, footnoteSlides = createFootnoteSlides(
            prs, slideNumber, footnoteDefinitions
        )

        footnotesPerPage = md2pptx.globals.processingOptions.getCurrentOption("footnotesPerPage")
        # Fix up any footnote slide hyperlinks
        footnoteNumber = -1
        for footnoteRun in footnoteRunsDictionary.keys():
            footnoteNumber += 1
            run = footnoteRunsDictionary[footnoteRun]

            footnoteSlideNumber = int(footnoteNumber / footnotesPerPage)
            createRunHyperlinkOrTooltip(run, footnoteSlides[footnoteSlideNumber], "")

    ######################################################################################
    #                                                                                    #
    # Add a dictionary slide - if there were any abbr elements encountered               #
    #                                                                                    #
    ######################################################################################
    if len(abbrevDictionary) > 0:
        glossaryTermsPerPage = md2pptx.globals.processingOptions.getCurrentOption("glossaryTermsPerPage")
        slideNumber, glossarySlides = createGlossarySlides(
            prs, slideNumber, abbrevDictionary
        )
        # Fix up internal glossary hyperlinks
        abbrevNumber = -1
        for abbreviation in sorted(abbrevRunsDictionary.keys()):
            abbrevNumber += 1
            runs = abbrevRunsDictionary[abbreviation]
            for run in runs:
                # Add tooltip for glossary definition
                glossarySlideNumber = int(abbrevNumber / glossaryTermsPerPage)
                createRunHyperlinkOrTooltip(
                    run, glossarySlides[glossarySlideNumber], abbrevDictionary[abbreviation]
                )

    ######################################################################################
    #                                                                                    #
    # Add final slide - or more than one - with any Taskpaper tasks in                   #
    #                                                                                    #
    ######################################################################################
    taskSlides = md2pptx.globals.processingOptions.getCurrentOption("taskSlides")

    if (len(tasks) > 0) & (taskSlides != "none"):
        # Turn tasks into a table slide

        # Might need to winnow slides
        if taskSlides != "all":
            complete = []
            incomplete = []
            for task in tasks:
                sNum, taskText, dueDate, tags, done = task
                if done == "":
                    incomplete.append(task)
                else:
                    complete.append(task)

            if (taskSlides == "separate") & (len(tasks) > 0):
                want_task_slides = True
            elif (taskSlides == "remaining") & (len(incomplete) > 0):
                want_task_slides = True
            elif (taskSlides == "done") & (len(complete) > 0):
                want_task_slides = True
            else:
                want_task_slides = False
        else:
            want_task_slides = True

        if want_task_slides:
            if taskSlides != "separate":
                createTaskSlides(prs, slideNumber, tasks, "Tasks")
            else:
                createTaskSlides(prs, slideNumber, complete, "Completed Tasks")
                createTaskSlides(prs, slideNumber, incomplete, "Incomplete Tasks")

    ######################################################################################
    #                                                                                    #
    # Make any TOC / Section-related links                                               #
    #                                                                                    #
    ######################################################################################
    if md2pptx.globals.processingOptions.getCurrentOption("tocLinks"):
        # Linkify section items
        for run in TOCruns:
            createRunHyperlinkOrTooltip(run, SectionSlides[run.text])

    if md2pptx.globals.processingOptions.getCurrentOption("sectionArrows"):
        # Add navigation arrows between section slides
        sectionArrowsColour = md2pptx.globals.processingOptions.getCurrentOption("sectionArrowsColour")

        buttonTop = prs.slide_height - Inches(2 / 3)
        forwShape = None

        previousSection = None
        for sectionNumber, sectionSlide in enumerate(SectionSlides):
            slide = SectionSlides[sectionSlide]
            if sectionNumber == 0:
                TOCslide = slide

            if forwShape != None:
                createShapeHyperlinkAndTooltip(forwShape, slide, "Next Section")

            buttonShapes = []
            if sectionNumber > 1:
                # Need backwards arrow
                backShape = slide.shapes.add_shape(
                    MSO_SHAPE.ACTION_BUTTON_BACK_OR_PREVIOUS,
                    prs.slide_width / 2 - Inches(2 / 3),
                    buttonTop,
                    Inches(1 / 3),
                    Inches(1 / 3),
                )

                createShapeHyperlinkAndTooltip(backShape, previousSlide, "Previous Section")

                buttonShapes.append(backShape)

            # Always need home arrow - except for TOC
            if sectionNumber > 0:
                homeShape = slide.shapes.add_shape(
                    MSO_SHAPE.ACTION_BUTTON_HOME,
                    prs.slide_width / 2 - Inches(1 / 6),
                    buttonTop,
                    Inches(1 / 3),
                    Inches(1 / 3),
                )

                createShapeHyperlinkAndTooltip(homeShape, TOCslide, "Table Of Contents")

                buttonShapes.append(homeShape)

            if (sectionNumber < len(SectionSlides) - 1) & (sectionNumber > 0):
                # Need forwards
                forwShape = slide.shapes.add_shape(
                    MSO_SHAPE.ACTION_BUTTON_FORWARD_OR_NEXT,
                    prs.slide_width / 2 + Inches(1 / 3),
                    buttonTop,
                    Inches(1 / 3),
                    Inches(1 / 3),
                )

                buttonShapes.append(forwShape)

            else:
                forwShape = None

            # Fix background colour of the buttons
            if sectionArrowsColour != "":
                for buttonShape in buttonShapes:
                    buttonShape.fill.solid()
                    buttonShape.fill.fore_color.rgb = RGBColor.from_string(
                        sectionArrowsColour
                    )

            previousSlide = slide


    ######################################################################################
    #                                                                                    #
    # Fix up any internal links                                                          #
    #                                                                                    #
    ######################################################################################
    xrefCheck_errors = False
    for href in md2pptx.globals.href_runs.keys():
        run = md2pptx.globals.href_runs[href]
        if href in slideHrefs.keys():
            createRunHyperlinkOrTooltip(
                run, prs.slides[slideHrefs[href] - 2 + templateSlideCount], ""
            )
        else:
            # No id defined with that name
            if not xrefCheck_errors:
                # First time in this run a cross reference error occurred
                xrefCheck_errors = True
                print("\nHyperlink Reference Errors")
                print("--------------------------")

            print(
                "Slide "
                + str(prs.slides.index(SlideFromRun(run)) + 1 - templateSlideCount)
                + f": '{href}'"
            )

    # Each picture appears in pictures
    # There's a corresponding entry in picture_Hrefs
    # There's a corresponding entry in picture_tooltips

    # fix up any clickable picture links
    for (picture, href, tooltip) in pictureInfos:
        # Pick up link target - if any
        if href == None:
            target = None
        else:
            if href[1:] in slideHrefs.keys():
                # Is an internal link
                target = prs.slides[slideHrefs[href[1:]] - 2 + templateSlideCount]
            else:
                # Is an external Link
                target = href

        createPictureHyperlinkOrTooltip(picture, target, tooltip)

    if templateSlideCount > 0:
        createProcessingSummarySlide(prs, metadata)

    try:
        if md2pptx.globals.processingOptions.getCurrentOption("deletefirstslide"):
            # Remove first slide
            deleteSlide(prs, 0)

    except:
        pass

    if md2pptx.globals.processingOptions.getCurrentOption("SectionsExpand"):
        createExpandingSections(prs)

    # Maybe call an exit before the presentation is saved
    onPresentationBeforeSave = md2pptx.globals.processingOptions.getCurrentOption("onPresentationBeforeSave")
    if onPresentationBeforeSave != "":
        exec(open(onPresentationBeforeSave).read())

    prs.save(output_filename)

    # Maybe call an exit after the presentation is saved
    onPresentationAfterSave = md2pptx.globals.processingOptions.getCurrentOption("onPresentationAfterSave")
    if onPresentationAfterSave != "":
        exec(open(onPresentationAfterSave).read())


    elapsed_time = time.time() - start_time

    print(
        "\nProcessing complete.\nElapsed time: "
        + str(int(1000 * elapsed_time) / 1000)
        + "s"
    )

    # Run a script against every slide
    script = "if slide.slideInfo is not None:\n  print(slide.slideInfo)"
    script = ""
    for slide in prs.slides:
        exec(script)


    sys.exit()
